"""Tests for OOXML font embedding (core/font_embed.py)."""

from __future__ import annotations

import zipfile
from pathlib import Path

import pytest

from core import font_embed

_FONTS_DIR = Path(__file__).resolve().parents[1] / "assets" / "fonts"


def _minimal_pptx(path: Path) -> None:
    """Write a tiny one-slide .pptx via python-pptx for the embed round-trip."""
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Hello"
    prs.save(str(path))


def test_font_file_for_resolves_family() -> None:
    """A family name maps to its space-stripped TTF, and unknown families resolve to None."""
    assert font_embed.font_file_for("Space Grotesk", _FONTS_DIR) == _FONTS_DIR / "SpaceGrotesk.ttf"
    assert font_embed.font_file_for("Not A Real Font", _FONTS_DIR) is None


def test_embed_fonts_injects_ooxml_parts(tmp_path: Path) -> None:
    """Embedding adds the font parts, the embeddedFontLst, the rel, and the content-type default."""
    pytest.importorskip("pptx")
    deck = tmp_path / "deck.pptx"
    _minimal_pptx(deck)

    ok = font_embed.embed_fonts(deck, ["Space Grotesk", "Inter", "Space Mono"], _FONTS_DIR)
    assert ok is True

    with zipfile.ZipFile(deck) as archive:
        names = set(archive.namelist())
        presentation = archive.read("ppt/presentation.xml").decode("utf-8")
        content_types = archive.read("[Content_Types].xml").decode("utf-8")
        rels = archive.read("ppt/_rels/presentation.xml.rels").decode("utf-8")

    assert {"ppt/fonts/font1.fntdata", "ppt/fonts/font3.fntdata"} <= names
    assert "embeddedFontLst" in presentation and 'embedTrueTypeFonts="1"' in presentation
    assert "fntdata" in content_types
    assert "relationships/font" in rels

    # The patched package must still be a valid presentation python-pptx can reopen.
    import pptx

    assert len(pptx.Presentation(str(deck)).slides) == 1


def test_embed_fonts_dedupes_and_skips_unknown(tmp_path: Path) -> None:
    """Duplicate families collapse to one part; families without a TTF are skipped."""
    pytest.importorskip("pptx")
    deck = tmp_path / "deck.pptx"
    _minimal_pptx(deck)

    ok = font_embed.embed_fonts(deck, ["Inter", "Inter", "Not A Real Font"], _FONTS_DIR)
    assert ok is True
    with zipfile.ZipFile(deck) as archive:
        font_parts = [n for n in archive.namelist() if n.startswith("ppt/fonts/")]
    assert font_parts == ["ppt/fonts/font1.fntdata"]


def test_embed_fonts_no_resolvable_families_is_noop(tmp_path: Path) -> None:
    """With no resolvable family the deck is returned untouched (fail-open)."""
    pytest.importorskip("pptx")
    deck = tmp_path / "deck.pptx"
    _minimal_pptx(deck)
    before = deck.read_bytes()

    assert font_embed.embed_fonts(deck, ["Not A Real Font"], _FONTS_DIR) is False
    assert deck.read_bytes() == before
