"""Tests for the PowerPoint (.pptx) reader (Dimensions / Phase 6).

Builds a real .pptx with python-pptx (title, textbox, table, speaker notes) in a temp dir,
then asserts the reader extracts all of them. No network, no LLM — pure local round-trip.
"""

from __future__ import annotations

from pathlib import Path

import pytest

from core.pptx_reader import PptxReadError, read_pptx


def _make_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # "Title Only"
    slide.shapes.title.text = "Quarterly Review"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(1))
    textbox.text_frame.text = "Revenue grew twelve percent"

    table = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(5), Inches(1)).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "EUR 4.2M"

    slide.notes_slide.notes_text_frame.text = "Speaker note: figures from the Q2 board pack"
    presentation.save(str(path))


def test_read_pptx_extracts_text_tables_and_notes(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.pptx"
    _make_pptx(deck_path)

    result = read_pptx(deck_path)

    assert result.doc_type == "pptx"
    assert result.extraction_method == "python-pptx"
    assert result.page_count == 1
    assert "Quarterly Review" in result.text
    assert "Revenue grew twelve percent" in result.text
    assert "Metric\tValue" in result.text
    assert "Revenue\tEUR 4.2M" in result.text
    assert "Speaker note: figures from the Q2 board pack" in result.text


def test_read_pptx_missing_file_raises(tmp_path: Path) -> None:
    with pytest.raises(FileNotFoundError):
        read_pptx(tmp_path / "missing.pptx")


def test_read_pptx_wrong_suffix_raises(tmp_path: Path) -> None:
    not_pptx = tmp_path / "deck.txt"
    not_pptx.write_text("plain text", encoding="utf-8")
    with pytest.raises(PptxReadError):
        read_pptx(not_pptx)
