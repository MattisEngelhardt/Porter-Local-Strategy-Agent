"""Integration tests for the composable block library (Block 2.2).

Each block is rendered onto a real python-pptx slide through the live ``_DeckRenderer`` (which
structurally satisfies the :class:`~core.blocks.Surface` protocol), then the resulting shapes are
inspected. This exercises both the blocks and the inch-based surface primitives the renderer adds.
"""

from __future__ import annotations

import pytest

from core import blocks
from core.blocks import BlockTheme
from core.config import AppConfig
from core.layout import Region
from models.task import Language
from models.visuals import ChartSeries, ChartSpec, ChartType

pytest.importorskip("pptx")


def _surface() -> object:
    from core.exporter import _DeckRenderer

    return _DeckRenderer(AppConfig(), Language.EN)


def _theme(renderer: object) -> BlockTheme:
    colors = renderer.colors  # type: ignore[attr-defined]
    return BlockTheme(
        colors=colors,
        fonts=renderer.fonts,  # type: ignore[attr-defined]
        editorial=True,
        fg=colors.ink,
        on_dark=False,
        spot=colors.coral,
        muted=colors.charcoal,
        accent=colors.accent_cyan,
    )


def _texts(slide: object) -> list[str]:
    out: list[str] = []
    for shape in slide.shapes:  # type: ignore[attr-defined]
        if shape.has_text_frame:
            out.append(shape.text_frame.text)
        elif shape.has_table:
            for row in shape.table.rows:
                out.extend(cell.text for cell in row.cells)
    return out


def _pictures(slide: object) -> int:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    return sum(
        1
        for shape in slide.shapes  # type: ignore[attr-defined]
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )


def test_block_kinds_cover_the_library() -> None:
    for kind in ("headline", "bullets", "cards", "stat_tiles", "flow", "matrix", "table", "chart"):
        assert kind in blocks.BLOCK_KINDS


def test_headline_renders_full_text_and_kicker() -> None:
    r = _surface()
    slide = r._new()  # type: ignore[attr-defined]
    region = Region(0.9, 0.55, 11.5, 1.06)
    blocks.render(
        "headline",
        r,  # type: ignore[arg-type]
        slide,
        region,
        {"text": "Funding leaders pull ahead", "kicker": "Market", "accent": "#4DACC7"},
        _theme(r),
    )
    joined = " ".join(_texts(slide))
    assert "Funding leaders pull ahead" in joined  # multi-run line is not truncated/split
    assert "MARKET" in joined  # kicker rendered, upper-cased


def test_bullets_render_every_item() -> None:
    r = _surface()
    slide = r._new()  # type: ignore[attr-defined]
    region = Region(0.72, 1.82, 11.75, 4.9)
    items = ["Signal one", "Signal two", "Signal three"]
    blocks.render("bullets", r, slide, region, {"items": items}, _theme(r))  # type: ignore[arg-type]
    joined = " ".join(_texts(slide))
    for item in items:
        assert item in joined


def test_flow_keeps_long_node_text_without_hard_truncation() -> None:
    r = _surface()
    slide = r._new()  # type: ignore[attr-defined]
    region = Region(0.72, 2.55, 11.75, 2.4)
    long_label = "Establish a marquee industrial partnership within twelve months"
    blocks.render(
        "flow",
        r,  # type: ignore[arg-type]
        slide,
        region,
        {"nodes": ["Assess the field", long_label]},
        _theme(r),
    )
    joined = " ".join(_texts(slide))
    # the full multi-word label survives (wrapped/auto-fit), not sliced to a 7-word fragment
    assert long_label in joined


def test_table_minimal_with_emphasis_column_renders_cells() -> None:
    r = _surface()
    slide = r._new()  # type: ignore[attr-defined]
    region = Region(0.72, 1.82, 11.75, 4.0)
    rows = [["Metric", "Neura", "Rival"], ["Funding", "935", "130"], ["Robots", "12", "8"]]
    blocks.render(
        "table",
        r,  # type: ignore[arg-type]
        slide,
        region,
        {"rows": rows, "style": "minimal", "emphasis_col": 1},
        _theme(r),
    )
    joined = " ".join(_texts(slide))
    assert "Neura" in joined and "935" in joined and "Robots" in joined


def test_chart_block_places_a_picture_when_spec_present() -> None:
    pytest.importorskip("matplotlib")
    r = _surface()
    slide = r._new()  # type: ignore[attr-defined]
    region = Region(0.72, 1.82, 11.75, 3.6)
    spec = ChartSpec(
        chart_type=ChartType.COLUMN,
        categories=["Neura", "Rival"],
        series=[ChartSeries(name="Funding", values=[935.0, 130.0])],
        unit="m",
    )
    blocks.render("chart", r, slide, region, {"spec": spec}, _theme(r))  # type: ignore[arg-type]
    assert _pictures(slide) == 1


def test_chart_block_noops_without_spec() -> None:
    r = _surface()
    slide = r._new()  # type: ignore[attr-defined]
    region = Region(0.72, 1.82, 11.75, 3.6)
    blocks.render("chart", r, slide, region, {}, _theme(r))  # type: ignore[arg-type]
    assert _pictures(slide) == 0


def test_unknown_kind_and_block_error_are_noops() -> None:
    r = _surface()
    slide = r._new()  # type: ignore[attr-defined]
    region = Region(1.0, 1.0, 4.0, 2.0)
    # unknown kind → no-op, no raise
    blocks.render("does_not_exist", r, slide, region, {}, _theme(r))  # type: ignore[arg-type]
    # a block fed a wrong-typed param fails open (the dispatcher swallows it)
    blocks.render("stat_tiles", r, slide, region, {"tiles": "not a list"}, _theme(r))  # type: ignore[arg-type]
    assert isinstance(_texts(slide), list)
