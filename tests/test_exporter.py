"""Tests for output rendering (core/exporter.py): PPTX (live) + PDF (fail-fast without GTK)."""

from __future__ import annotations

from pathlib import Path

import pytest

from core.artifact_framework import framework_marker
from core.config import AppConfig, StyleConfig
from core.exporter import (
    PdfBuildError,
    brief_template_for,
    build_brief_pdf,
    build_deck,
    build_management_deck,
    build_management_pdf,
    render_brief_html,
)
from models.deck import DeckStructure, SlideContent, SlideType
from models.research import Confidence, Finding, ResearchReport, WorkerFindings
from models.synthesis import AnalysisOutput, Section, SourceRef
from models.task import Language, TaskType
from models.visuals import ChartSeries, ChartSpec, ChartType


def _presentation_text(prs) -> str:
    """Collect text from text boxes and table cells in a rendered PPTX."""
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    chunks.extend(cell.text for cell in row.cells)
    return " ".join(chunks)


def _analysis() -> AnalysisOutput:
    return AnalysisOutput(
        title="Q2 Board Update",
        language=Language.EN,
        bottom_line="Runway is 9 months. Approve the bridge round now.",
        sections=[
            Section(
                heading="Cash runway shortens to 9 months", body="Cash fell to 10.8M. Burn 1.2M/mo."
            ),
            Section(heading="Revenue up 35%", body="Revenue rose to 4.2M on industrial pilots."),
        ],
        sources=[SourceRef(url="q2_financials.xlsx", title="Q2 figures")],
    )


def test_build_management_deck_creates_pptx(tmp_path: Path) -> None:
    """A management .pptx is rendered through the mandatory artifact framework."""
    pptx = pytest.importorskip("pptx")  # python-pptx is a declared dependency
    path = build_management_deck(_analysis(), AppConfig(), tmp_path, Language.EN)
    assert path.exists() and path.suffix == ".pptx"

    prs = pptx.Presentation(str(path))
    # title + executive summary + evidence + 2 sections + recommendation + sources
    assert len(prs.slides) == 7
    title_texts = [shape.text_frame.text for shape in prs.slides[0].shapes if shape.has_text_frame]
    assert any("Q2 Board Update" in text for text in title_texts)
    all_text = _presentation_text(prs)
    assert framework_marker() in all_text
    assert "The source base shows" in all_text
    assert "The next management move follows from the bottom line" in all_text
    assert "Cash runway shortens to 9 months" in all_text


def test_build_management_pdf_failfast_or_renders(tmp_path: Path) -> None:
    """PDF renders when WeasyPrint+GTK are available; otherwise fails fast with fix instructions."""
    try:
        import weasyprint  # noqa: F401

        available = True
    except Exception:  # ImportError or OSError (missing GTK runtime)
        available = False

    if available:
        path = build_management_pdf(_analysis(), AppConfig(), tmp_path, Language.EN)
        assert path.exists() and path.suffix == ".pdf"
    else:
        with pytest.raises(PdfBuildError) as excinfo:
            build_management_pdf(_analysis(), AppConfig(), tmp_path, Language.EN)
        message = str(excinfo.value).lower()
        assert "gtk" in message or "weasyprint" in message


# ----------------------------------------------------------------- brief HTML (pure)
def test_brief_template_routing() -> None:
    """Each task type maps to its SPEC §10 brief template (T-1..T-6)."""
    assert brief_template_for(TaskType.COMPETITOR_ANALYSIS) == "competitor_brief.md.j2"
    assert brief_template_for(TaskType.TARGET_SCREENING) == "decision_brief.md.j2"
    assert brief_template_for(TaskType.MARKET_ANALYSIS) == "market_overview.md.j2"
    assert brief_template_for(TaskType.BOARD_PREP) == "board_update.md.j2"
    assert brief_template_for(TaskType.DOCUMENT_SYNTHESIS) == "document_synthesis.md.j2"
    assert brief_template_for(TaskType.INDUSTRY_NEWS) == "adhoc_brief.md.j2"


def test_render_brief_html_structure_and_bullets() -> None:
    """The HTML leads with the bottom line, renders sections, and converts bullet lines."""
    analysis = AnalysisOutput(
        title="1X — Competitive Brief",
        language=Language.EN,
        bottom_line="1X is well funded; Neura must differentiate on cognition.",
        sections=[
            Section(heading="Three rivals are closing in", body="- 1X: $100M\n- Figure: $675M")
        ],
        sources=[SourceRef(url="https://reuters.com/x", title="round")],
    )
    html = render_brief_html(analysis, AppConfig(), task_type=TaskType.COMPETITOR_ANALYSIS)
    assert '<span class="hl">1X</span>' in html  # two-tone headline splits the key token (DNA 7)
    assert "Competitive Brief" in html
    assert "1X is well funded" in html
    assert "Three rivals are closing in" in html
    assert "<li>1X: $100M</li>" in html and "<li>Figure: $675M</li>" in html
    assert "Executive Summary" in html  # T-1 bottom-line label (EN)
    assert framework_marker() in html
    assert "artifact-ribbon" in html
    assert "Evidence anchors" in html
    assert "Focus" in html and "Proof" in html and "Sources" in html and "Status" in html
    assert "https://reuters.com/x" in html


def test_render_brief_html_is_bilingual() -> None:
    """A German analysis renders German labels; an English one renders English labels."""
    base = dict(title="T", bottom_line="b", sections=[Section(heading="h", body="x")], sources=[])
    de = render_brief_html(
        AnalysisOutput(language=Language.DE, **base),  # type: ignore[arg-type]
        AppConfig(),
        task_type=TaskType.TARGET_SCREENING,
    )
    en = render_brief_html(
        AnalysisOutput(language=Language.EN, **base),  # type: ignore[arg-type]
        AppConfig(),
        task_type=TaskType.TARGET_SCREENING,
    )
    assert "Empfehlung" in de and "Entscheidungsanalyse" in de
    assert "Recommendation" in en and "Decision Analysis" in en


def test_render_brief_html_escapes_markup() -> None:
    """User/LLM text is HTML-escaped (no injection through the analysis)."""
    analysis = AnalysisOutput(
        title="A <script>alert(1)</script>",
        language=Language.EN,
        bottom_line="b & c < d",
        sections=[],
        sources=[],
    )
    html = render_brief_html(analysis, AppConfig(), task_type=TaskType.ADHOC)
    assert "<script>alert(1)</script>" not in html
    assert "&lt;script&gt;" in html
    assert "b &amp; c &lt; d" in html


# ----------------------------------------------------------------- brief (Editorial design system)
def test_render_brief_html_wires_editorial_multifont_and_canvas() -> None:
    """The brief renders on the cream canvas with the serif/grotesk/body/mono CSS stacks wired."""
    style = StyleConfig()
    colors = AppConfig().output.colors
    html = render_brief_html(_analysis(), AppConfig(), task_type=TaskType.BOARD_PREP)
    # cream editorial canvas + ink text (not the old white/blue Neura sheet)
    assert f"background: {colors.paper}" in html
    # multi-font system: every role's primary family is present in the stacks (fallbacks kept)
    assert style.serif_font in html  # serif display headlines (PDF)
    assert style.grotesk_font in html  # grotesk metric numerals
    assert style.body_font in html  # body text
    assert style.mono_font in html  # mono micro-labels / telemetry


def test_render_brief_html_editorial_has_cover_band_restrained_omits_it() -> None:
    """Editorial intensity lays a luminous gradient cover band; restrained drops it (board-safe)."""
    editorial = render_brief_html(_analysis(), AppConfig(), task_type=TaskType.BOARD_PREP)
    assert 'class="cover-band"' in editorial and "linearGradient" in editorial

    config = AppConfig()
    config.output.style.intensity = "restrained"
    restrained = render_brief_html(_analysis(), config, task_type=TaskType.BOARD_PREP)
    assert 'class="cover-band"' not in restrained  # the gradient element is gone
    assert "linearGradient" not in restrained  # no gradient anywhere in the restrained sheet


def test_render_brief_html_two_tone_headline_highlights_number() -> None:
    """A numeric headline token is wrapped in the coral spot-color span (DNA 7)."""
    analysis = AnalysisOutput(
        title="Revenue grew 40% in Q2",
        language=Language.EN,
        bottom_line="Momentum is real.",
        sections=[Section(heading="Why", body="Industrial pilots converted.")],
        sources=[],
    )
    html = render_brief_html(analysis, AppConfig(), task_type=TaskType.ADHOC)
    assert '<span class="hl">40%</span>' in html


def test_render_brief_html_renders_telemetry_only_with_report() -> None:
    """Source-grounded telemetry chips appear only when a research report is threaded in (DNA 6)."""
    with_report = render_brief_html(
        _analysis(),
        AppConfig(),
        task_type=TaskType.BOARD_PREP,
        research_report=_research_report(),
    )
    assert "SOURCES 22" in with_report and "WORKERS 3" in with_report

    without = render_brief_html(_analysis(), AppConfig(), task_type=TaskType.BOARD_PREP)
    assert "SOURCES 22" not in without  # never invents telemetry


def test_render_brief_html_embeds_section_chart_svg() -> None:
    """A section carrying a ``.visual`` embeds a hand-built inline SVG chart (locked PDF tech)."""
    analysis = AnalysisOutput(
        title="Funding leaders pull ahead",
        language=Language.EN,
        bottom_line="Rivals are better capitalized.",
        sections=[
            Section(
                heading="Funding by company",
                body="Figure leads at 39M.",
                visual=_chart_spec(),
            )
        ],
        sources=[],
    )
    html = render_brief_html(analysis, AppConfig(), task_type=TaskType.MARKET_ANALYSIS)
    assert 'figure class="chart"' in html
    assert "<svg" in html and "Apptronik" in html  # a category label from the spec


def test_render_brief_html_chart_budget_caps_embeds() -> None:
    """No more than ``style.max_charts_per_brief`` section charts embed (rest render text-only)."""
    config = AppConfig()
    config.output.style.max_charts_per_brief = 1
    analysis = AnalysisOutput(
        title="Charts everywhere",
        language=Language.EN,
        bottom_line="Many numbers.",
        sections=[
            Section(heading=f"Section {i}", body="body", visual=_chart_spec()) for i in range(3)
        ],
        sources=[],
    )
    html = render_brief_html(analysis, config, task_type=TaskType.MARKET_ANALYSIS)
    assert html.count('figure class="chart"') == 1


def test_build_brief_pdf_failfast_or_renders(tmp_path: Path) -> None:
    """build_brief_pdf renders a .pdf when WeasyPrint+GTK are present, else fails fast."""
    try:
        import weasyprint  # noqa: F401

        available = True
    except Exception:  # ImportError or OSError (missing/!broken GTK runtime)
        available = False

    if available:
        path = build_brief_pdf(
            _analysis(), AppConfig(), tmp_path, task_type=TaskType.COMPETITOR_ANALYSIS
        )
        assert path.exists() and path.suffix == ".pdf"
    else:
        with pytest.raises(PdfBuildError):
            build_brief_pdf(_analysis(), AppConfig(), tmp_path, task_type=TaskType.ADHOC)


# ----------------------------------------------------------------- deck (all 10 slide types)
def _full_deck() -> DeckStructure:
    return DeckStructure(
        title="Neura Q2 Board",
        language=Language.EN,
        slides=[
            SlideContent(
                slide_type=SlideType.TITLE, headline="Neura Board Update Q2", body="Board"
            ),
            SlideContent(
                slide_type=SlideType.EXECUTIVE_SUMMARY,
                headline="Runway is 9 months",
                bullets=["Cash 10.8M", "Burn 1.2M/mo"],
            ),
            SlideContent(slide_type=SlideType.MARKET_LANDSCAPE, headline="Rivals closing in"),
            SlideContent(slide_type=SlideType.COMPANY_DEEP_DIVE, headline="Figure scaling"),
            SlideContent(slide_type=SlideType.FINANCIAL_OVERVIEW, headline="Funding up"),
            SlideContent(
                slide_type=SlideType.COMPETITIVE_COMPARISON,
                headline="Neura leads on cognition",
                table=[["Company", "Funding"], ["Neura", "120M"], ["Figure", "675M"]],
            ),
            SlideContent(slide_type=SlideType.STRATEGIC_SIGNALS, headline="Hiring signals push"),
            SlideContent(
                slide_type=SlideType.SWOT,
                headline="Strong tech, thin capital",
                table=[
                    ["Strengths", "Cognitive AI; Bosch"],
                    ["Weaknesses", "Less capital"],
                    ["Opportunities", "EU industrial"],
                    ["Threats", "US scale"],
                ],
            ),
            SlideContent(
                slide_type=SlideType.RECOMMENDATION,
                headline="Approve the bridge round now",
                body="GO — raise €80M bridge",
                bullets=["Extends runway 18mo"],
            ),
            SlideContent(
                slide_type=SlideType.APPENDIX, headline="Sources", bullets=["reuters.com/x"]
            ),
        ],
    )


def test_build_deck_all_slide_types_with_logo(tmp_path: Path) -> None:
    """build_deck renders all 10 slide types, one slide each, with the logo bottom-right."""
    pptx = pytest.importorskip("pptx")
    path = build_deck(_full_deck(), AppConfig(), tmp_path)
    assert path.exists() and path.suffix == ".pptx"

    prs = pptx.Presentation(str(path))
    assert len(prs.slides) == 10
    # Logo (a picture, shape_type == 13) appears on every slide (SPEC §11).
    for slide in prs.slides:
        assert any(shape.shape_type == 13 for shape in slide.shapes)
    # The "so what" recommendation headline and its decision callout are present.
    all_text = _presentation_text(prs)
    assert framework_marker() in all_text
    assert "Approve the bridge round now" in all_text
    assert "GO — raise" in all_text
    # The comparison table cells made it into the deck.
    assert "Cognitive AI" in all_text  # SWOT quadrant content


def test_build_deck_with_analysis_applies_prerender_framework(tmp_path: Path) -> None:
    """build_deck inserts the evidence/recommendation/source frame before rendering."""
    pptx = pytest.importorskip("pptx")
    deck = DeckStructure(
        title="Thin Draft",
        language=Language.EN,
        slides=[SlideContent(slide_type=SlideType.TITLE, headline="Thin Draft")],
    )
    path = build_deck(deck, AppConfig(), tmp_path, analysis=_analysis())

    prs = pptx.Presentation(str(path))
    all_text = _presentation_text(prs)
    assert framework_marker() in all_text
    assert "Executive Summary" in all_text
    assert "The source base shows" in all_text
    assert "The next management move follows from the bottom line" in all_text
    assert "q2_financials.xlsx" in all_text


def test_build_deck_no_logo_when_disabled(tmp_path: Path) -> None:
    """With include_logo off, no picture is added (the bottom-right logo is config-gated)."""
    pptx = pytest.importorskip("pptx")
    config = AppConfig()
    config.output.include_logo = False
    path = build_deck(_full_deck(), config, tmp_path)
    prs = pptx.Presentation(str(path))
    assert all(all(shape.shape_type != 13 for shape in slide.shapes) for slide in prs.slides)


# ----------------------------------------------------------------- deck (Editorial visual engine)
def _chart_spec() -> ChartSpec:
    return ChartSpec(
        chart_type=ChartType.COLUMN,
        categories=["Figure", "1X", "Apptronik"],
        series=[ChartSeries(name="Funding (EUR m)", values=[39.0, 12.0, 7.0])],
        caption="Funding by company",
        unit="m",
    )


def _research_report() -> ResearchReport:
    return ResearchReport(
        worker_findings=[
            WorkerFindings(
                sub_topic="funding",
                findings=[Finding(claim="Raised $55M", date="2025-01", confidence=Confidence.HIGH)],
            )
        ],
        workers_used=3,
        sources_evaluated=22,
    )


def _all_runs(prs: object) -> list[object]:
    """Every text run across the deck (for font/color assertions)."""
    runs: list[object] = []
    for slide in prs.slides:  # type: ignore[attr-defined]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    runs.extend(para.runs)
    return runs


def _gradient_count(prs: object) -> int:
    """Count shapes painted with a gradient fill (the editorial luminous depth)."""
    from pptx.enum.dml import MSO_FILL

    count = 0
    for slide in prs.slides:  # type: ignore[attr-defined]
        for shape in slide.shapes:
            try:
                if shape.fill.type == MSO_FILL.GRADIENT:
                    count += 1
            except (AttributeError, TypeError, ValueError):
                continue  # pictures / tables have no addressable fill
    return count


def test_build_deck_renders_native_chart_from_visual(tmp_path: Path) -> None:
    """A slide carrying a ``.visual`` renders a native (editable) python-pptx chart."""
    pptx = pytest.importorskip("pptx")
    deck = DeckStructure(
        title="Charted",
        language=Language.EN,
        slides=[
            SlideContent(slide_type=SlideType.TITLE, headline="Charted"),
            SlideContent(
                slide_type=SlideType.FINANCIAL_OVERVIEW,
                headline="Funding leaders pull ahead",
                visual=_chart_spec(),
            ),
        ],
    )
    path = build_deck(deck, AppConfig(), tmp_path)
    prs = pptx.Presentation(str(path))
    assert any(getattr(shape, "has_chart", False) for slide in prs.slides for shape in slide.shapes)


def test_build_deck_chart_budget_caps_visuals(tmp_path: Path) -> None:
    """No more than ``style.max_charts_per_deck`` native charts are rendered (rest fall back)."""
    pptx = pytest.importorskip("pptx")
    config = AppConfig()
    config.output.style.max_charts_per_deck = 1
    slides = [SlideContent(slide_type=SlideType.TITLE, headline="Many charts")]
    slides += [
        SlideContent(
            slide_type=SlideType.STRATEGIC_SIGNALS,
            headline=f"Signal {i}",
            visual=_chart_spec(),
        )
        for i in range(3)
    ]
    deck = DeckStructure(title="Capped", language=Language.EN, slides=slides)
    path = build_deck(deck, config, tmp_path)
    prs = pptx.Presentation(str(path))
    charts = sum(
        1 for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_chart", False)
    )
    assert charts == 1


def test_title_cover_uses_dark_editorial_canvas(tmp_path: Path) -> None:
    """The cover sits on the dramatic ``canvas_dark`` editorial canvas with a luminous gradient."""
    pptx = pytest.importorskip("pptx")
    colors = AppConfig().output.colors
    path = build_deck(_full_deck(), AppConfig(), tmp_path)
    prs = pptx.Presentation(str(path))
    background = prs.slides[0].background.fill.fore_color.rgb
    assert str(background) == colors.canvas_dark.lstrip("#").upper()
    assert _gradient_count(prs) >= 1  # editorial depth present on cover/divider


def test_restrained_intensity_skips_gradient(tmp_path: Path) -> None:
    """``restrained`` intensity keeps the structure but drops all gradient depth (board-safe)."""
    pptx = pytest.importorskip("pptx")
    config = AppConfig()
    config.output.style.intensity = "restrained"
    path = build_deck(_full_deck(), config, tmp_path / "restrained")
    prs = pptx.Presentation(str(path))
    assert _gradient_count(prs) == 0
    # but the cover is still the dark canvas (solid), so it stays a strong title moment
    canvas_dark = config.output.colors.canvas_dark.lstrip("#").upper()
    assert str(prs.slides[0].background.fill.fore_color.rgb) == canvas_dark


def test_build_deck_renders_telemetry_chips(tmp_path: Path) -> None:
    """Source-grounded telemetry chips render only when a research report is threaded in (DNA 6)."""
    pptx = pytest.importorskip("pptx")
    report = _research_report()
    path = build_deck(_full_deck(), AppConfig(), tmp_path / "with", research_report=report)
    prs = pptx.Presentation(str(path))
    all_text = _presentation_text(prs)
    assert "SOURCES 22" in all_text
    assert "WORKERS 3" in all_text

    # No report → no chips (never invents telemetry).
    path_without = build_deck(_full_deck(), AppConfig(), tmp_path / "without")
    prs_without = pptx.Presentation(str(path_without))
    assert "SOURCES 22" not in _presentation_text(prs_without)


def test_headline_two_tone_splits_numeric_token(tmp_path: Path) -> None:
    """A numeric headline is rendered in multiple runs with the number in the spot color (DNA 7)."""
    pptx = pytest.importorskip("pptx")
    deck = DeckStructure(
        title="Two tone",
        language=Language.EN,
        slides=[
            SlideContent(slide_type=SlideType.STRATEGIC_SIGNALS, headline="Revenue grew 40% in Q2")
        ],
    )
    path = build_deck(deck, AppConfig(), tmp_path)
    prs = pptx.Presentation(str(path))
    coral = AppConfig().output.colors.coral.lstrip("#").upper()
    headline_runs: list[object] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and "Revenue grew 40%" in shape.text_frame.text:
                headline_runs = shape.text_frame.paragraphs[0].runs
    assert len(headline_runs) >= 2  # split into base + accent runs
    assert any(str(run.font.color.rgb) == coral for run in headline_runs)  # token in spot color


def test_deck_shapes_are_valid_ooxml(tmp_path: Path) -> None:
    """Every rendered shape obeys the CT_ShapeProperties content model (PowerPoint-safe).

    PowerPoint *repairs and silently drops* a shape whose ``spPr`` has a duplicated
    single-occurrence child (e.g. two ``<a:effectLst>``) or out-of-order children — even though
    python-pptx happily reopens it. This pins the raw-DrawingML depth helpers (gradient/glow) so
    that class of corruption can never ship again. Covers the dark cover + recommendation divider
    (both carry the gradient wash + soft-edged glow).
    """
    pptx = pytest.importorskip("pptx")
    from pptx.oxml.ns import qn

    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    order = [
        "xfrm",
        "custGeom",
        "prstGeom",
        "noFill",
        "solidFill",
        "gradFill",
        "blipFill",
        "pattFill",
        "grpFill",
        "ln",
        "effectLst",
        "effectDag",
        "scene3d",
        "sp3d",
        "extLst",
    ]
    fills = {"noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill"}
    path = build_deck(_full_deck(), AppConfig(), tmp_path, research_report=_research_report())
    prs = pptx.Presentation(str(path))
    for slide in prs.slides:
        for shape in slide.shapes:
            sppr = shape._element.find(qn("p:spPr"))
            if sppr is None:
                continue
            tags = [
                child.tag.split("}")[-1]
                for child in sppr
                if child.tag.startswith(f"{{{ns}}}") and child.tag.split("}")[-1] in order
            ]
            for tag in set(tags):
                assert tags.count(tag) == 1, f"duplicate <a:{tag}> in {shape.name}: {tags}"
            assert sum(tags.count(f) for f in fills) <= 1, f"multiple fills in {shape.name}: {tags}"
            indices = [order.index(t) for t in tags]
            assert indices == sorted(indices), f"out-of-order spPr in {shape.name}: {tags}"


def test_deck_uses_editorial_multifont_system(tmp_path: Path) -> None:
    """Headlines use the grotesk display font, micro-labels the mono font, body the body font."""
    pptx = pytest.importorskip("pptx")
    style = StyleConfig()
    path = build_deck(_full_deck(), AppConfig(), tmp_path)
    prs = pptx.Presentation(str(path))
    fonts = {run.font.name for run in _all_runs(prs) if run.font.name}
    assert style.grotesk_font in fonts  # two-tone grotesk headlines
    assert style.mono_font in fonts  # frame label / page number / telemetry / tags
    assert style.body_font in fonts  # card body + table + bullets
    assert "Arial" not in fonts  # the old hardcoded bullet font is gone
