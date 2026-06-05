"""Architecture guards: the design/render engine is model- and provider-agnostic.

Because all design is deterministic code behind the typed content contract (``DeckStructure`` /
``SlideContent`` / ``ChartSpec`` …), swapping the LLM (4B → 12B → an API) changes only *content*,
never *design*. These tests lock that invariant: no rendering module may import the LLM client, and
the same ``DeckStructure`` must render to the same deck regardless of which model produced it.
"""

from __future__ import annotations

import ast
from pathlib import Path

import pytest

from core.config import AppConfig
from core.exporter import build_deck
from models.deck import DeckStructure, SlideContent, SlideType
from models.task import Language

_ROOT = Path(__file__).resolve().parents[1]

# The design/render engine — none of these may depend on the model/provider.
_DESIGN_MODULES = [
    "core/exporter.py",
    "core/charts_image.py",
    "core/font_embed.py",
    "core/deck_director.py",
    "core/design.py",
    "core/imagery.py",
    "core/typography.py",
    "core/visuals.py",
    "core/visual_selector.py",
    "core/diagrams.py",
    "core/layout.py",
]


@pytest.mark.parametrize("rel", _DESIGN_MODULES)
def test_design_modules_do_not_import_the_llm(rel: str) -> None:
    """A rendering/design module must not import ``llm`` — design stays model-agnostic."""
    tree = ast.parse((_ROOT / rel).read_text(encoding="utf-8"))
    imported: set[str] = set()
    for node in ast.walk(tree):
        if isinstance(node, ast.Import):
            imported.update(alias.name for alias in node.names)
        elif isinstance(node, ast.ImportFrom) and node.module:
            imported.add(node.module)
    offenders = {name for name in imported if name == "llm" or name.startswith("llm.")}
    assert not offenders, f"{rel} imports llm {offenders}; design must be model-agnostic"


def _texts(prs: object) -> list[str]:
    chunks: list[str] = []
    for slide in prs.slides:  # type: ignore[attr-defined]
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
    return chunks


def test_same_deckstructure_renders_identically(tmp_path: Path) -> None:
    """Same content → same deck: the render is a pure function of the DeckStructure."""
    pytest.importorskip("pptx")
    import pptx

    deck = DeckStructure(
        title="Model agnostic",
        language=Language.EN,
        slides=[
            SlideContent(slide_type=SlideType.TITLE, headline="A board update"),
            SlideContent(
                slide_type=SlideType.STRATEGIC_SIGNALS,
                headline="Momentum builds",
                bullets=["Signal one", "Signal two", "Signal three"],
            ),
        ],
    )
    config = AppConfig()
    first = build_deck(deck, config, tmp_path / "a")
    second = build_deck(deck, config, tmp_path / "b")
    assert _texts(pptx.Presentation(str(first))) == _texts(pptx.Presentation(str(second)))
