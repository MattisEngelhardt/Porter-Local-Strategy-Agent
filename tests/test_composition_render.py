"""Integration tests for the Block-2 composition render path (2.5).

Builds real decks through ``build_deck`` (now composer-driven) and asserts the new templates land:
the rebuilt flow keeps long node text, the decision slide is restrained Neura b/w with a derived
chip, the metric hero shows a compact numeral (not a wrapping unit word), and the legacy
``render(sc, plan)`` archetype path still renders when no composition is supplied (REQ-5 fallback).
"""

from __future__ import annotations

from pathlib import Path

import pytest

from core import deck_director
from core.config import AppConfig
from core.exporter import build_deck
from models.deck import DeckStructure, SlideContent, SlideType
from models.task import Language

pytest.importorskip("pptx")


def _texts(prs: object) -> list[str]:
    out: list[str] = []
    for slide in prs.slides:  # type: ignore[attr-defined]
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
    return out


def _no_imagery(tmp_path: Path) -> AppConfig:
    config = AppConfig()
    config.output.imagery_dir = str(tmp_path / "no-imagery")  # deterministic, no image pickups
    return config


def test_flow_renders_steps_without_truncation(tmp_path: Path) -> None:
    import pptx

    long_step = "Secure one marquee deployment partner within twelve months across the EU"
    deck = DeckStructure(
        title="Signals",
        language=Language.EN,
        slides=[
            SlideContent(slide_type=SlideType.TITLE, headline="Signals"),
            SlideContent(
                slide_type=SlideType.STRATEGIC_SIGNALS,
                headline="Three moves",
                bullets=["Assess the field", long_step, "Scale the cognitive stack"],
            ),
        ],
    )
    path = build_deck(deck, _no_imagery(tmp_path), tmp_path)
    joined = " ".join(_texts(pptx.Presentation(str(path))))
    assert long_step in joined  # the flow wraps/auto-fits — no 7-word hard truncation
    assert "→" in joined  # connected step cards


def test_decision_slide_is_restrained_neura_with_a_chip(tmp_path: Path) -> None:
    import pptx

    deck = DeckStructure(
        title="Decision",
        language=Language.EN,
        slides=[
            SlideContent(slide_type=SlideType.TITLE, headline="Decision"),
            SlideContent(
                slide_type=SlideType.RECOMMENDATION,
                headline="Accelerate — secure one marquee partner",
                body="Proceed now while the window is open.",
                bullets=["Sign a flagship pilot (COO, Q3)", "Raise a strategic round (CFO, H2)"],
            ),
        ],
    )
    config = _no_imagery(tmp_path)
    path = build_deck(deck, config, tmp_path)
    prs = pptx.Presentation(str(path))
    texts = _texts(prs)
    assert "GO" in texts  # the derived decision chip (from "Accelerate/secure/proceed")
    assert any("Accelerate" in t for t in texts)
    assert any("Sign a flagship pilot" in t for t in texts)
    # restrained near-white Neura canvas (not a saturated field) on the decision slide
    decision_slide = next(
        s
        for s in prs.slides
        if any(sh.has_text_frame and "Accelerate" in sh.text_frame.text for sh in s.shapes)
    )
    bg = str(decision_slide.background.fill.fore_color.rgb)
    assert bg == config.output.colors.cream_hi.lstrip("#").upper()


def test_metric_hero_token_is_a_compact_numeral(tmp_path: Path) -> None:
    import pptx

    deck = DeckStructure(
        title="Metrics",
        language=Language.EN,
        slides=[
            SlideContent(slide_type=SlideType.TITLE, headline="Metrics"),
            SlideContent(
                slide_type=SlideType.EXECUTIVE_SUMMARY,
                headline="The numbers",
                bullets=["Funding reached 675m", "Window is 12 months", "Margin hit 40%"],
            ),
        ],
    )
    path = build_deck(deck, _no_imagery(tmp_path), tmp_path)
    texts = _texts(pptx.Presentation(str(path)))
    # the giant numeral is the compact token ("12"), never the wrapping unit word ("12 months")
    assert "12" in texts
    assert "12 months" not in texts  # the unit word lives only in the supporting label
    assert "675m" in texts and "40%" in texts


def test_legacy_render_without_composition_still_works(tmp_path: Path) -> None:
    """``render(sc, plan)`` with no composition uses the archetype path — the ultimate fallback."""
    from core.exporter import _DeckRenderer

    sc = SlideContent(
        slide_type=SlideType.STRATEGIC_SIGNALS,
        headline="Momentum builds",
        bullets=["Signal one", "Signal two", "Signal three"],
    )
    renderer = _DeckRenderer(AppConfig(), Language.EN)
    plan = deck_director.plan_deck([sc])[0]
    before = len(renderer.prs.slides)
    renderer.render(sc, plan)  # comp omitted → legacy archetype path
    assert len(renderer.prs.slides) == before + 1
