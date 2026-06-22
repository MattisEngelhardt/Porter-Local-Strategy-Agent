"""PowerPoint (.pptx) reader (Dimensions / Phase 6): extract slide text, tables, and notes.

Porter already *writes* PPTX decks (:mod:`core.exporter`); this is the *reading* counterpart for
ingesting internal decks (board packs, strategy decks) into the Analyst/Builder dimensions.
Uses **python-pptx** (already a project dependency for writing, MIT-licensed, fully local) — so
reading PPTX adds no new dependency. Extraction only — no synthesis.

Each slide contributes its shape text and table cells, plus the speaker notes (where the real
argument behind a slide often lives). For high-fidelity layout on complex decks the optional
Docling adapter (:mod:`core.docling_reader`) is preferred; this is the always-available path.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from models.research import DocContent


class PptxReadError(Exception):
    """A .pptx document could not be read by python-pptx (fail fast, SPEC REQ-5)."""


# --------------------------------------------------------------- backend seam
def _shape_lines(shape: Any) -> list[str]:
    """Collect text from one shape: text frames (verbatim) and table cells (tab-separated)."""
    lines: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text.strip()
        if text:
            lines.append(text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                lines.append("\t".join(cells))
    return lines


def _extract_pptx(path: Path) -> tuple[str, int]:
    """Return (joined per-slide text incl. notes, slide count) via python-pptx."""
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - dependency is declared in requirements
        raise PptxReadError(
            "python-pptx is required to read PowerPoint (.pptx) files but is not installed.\n"
            "Fix: pip install python-pptx"
        ) from exc

    presentation = Presentation(str(path))
    parts: list[str] = []
    slide_count = 0

    for index, slide in enumerate(presentation.slides, start=1):
        slide_count += 1
        slide_lines: list[str] = []
        for shape in slide.shapes:
            slide_lines.extend(_shape_lines(shape))

        notes = ""
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()

        if slide_lines or notes:
            parts.append(f"--- Slide {index} ---")
            parts.extend(slide_lines)
            if notes:
                parts.append(f"[Notes] {notes}")

    return "\n".join(parts).strip(), slide_count


# --------------------------------------------------------------------- public API
def read_pptx(path: Path | str) -> DocContent:
    """Read a PowerPoint ``.pptx`` into :class:`DocContent` (slide text + tables + notes).

    Args:
        path: Path to a ``.pptx`` file.

    Returns:
        :class:`DocContent` with ``doc_type="pptx"``, ``page_count`` = slide count, and
        ``extraction_method="python-pptx"``.

    Raises:
        FileNotFoundError: If the path does not exist.
        PptxReadError: If the type is unsupported or the file is not a readable .pptx.
    """
    file_path = Path(path)
    if not file_path.is_file():
        raise FileNotFoundError(
            f"Document not found: {file_path.resolve()}\n"
            "Fix: check the path; it must point to an existing .pptx file."
        )

    if file_path.suffix.lower() != ".pptx":
        raise PptxReadError(
            f"Unsupported document type '{file_path.suffix}' for {file_path.name}. "
            "read_pptx handles .pptx only — convert legacy .ppt to .pptx first."
        )

    try:
        text, slide_count = _extract_pptx(file_path)
    except PptxReadError:
        raise
    except Exception as exc:  # python-pptx raises varied errors for corrupt/locked files
        raise PptxReadError(
            f"Could not read {file_path.name} as a PowerPoint document: {exc}\n"
            "Fix: ensure it is a valid, non-password-protected .pptx file."
        ) from exc

    return DocContent(
        source_path=file_path,
        doc_type="pptx",
        text=text,
        page_count=slide_count,
        extraction_method="python-pptx",
    )
