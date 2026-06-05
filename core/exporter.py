"""Output rendering (SPEC §7 ``exporter.py``): Neura-styled PDF briefs + PPTX decks.

Turns a structured :class:`~models.synthesis.AnalysisOutput` into delivery files with Neura
styling (colors from ``config.yaml``, logo bottom-right on decks / in the brief header):

* **PDF brief** — Jinja2 HTML templates (``templates/briefs/`` T-1..T-6, SPEC §10) → ``weasyprint``
  (the SPEC §6 PDF tool). :func:`render_brief_html` is pure/testable; :func:`build_brief_pdf` writes
  the PDF. On Windows WeasyPrint needs the GTK runtime — :func:`_ensure_gtk_dll_dir` forces a found
  GTK ``bin`` ahead of any incompatible ``libgobject`` on PATH, and if WeasyPrint still cannot be
  imported the call fails fast with exact install instructions (the renderer is correct and works
  the moment GTK is present — zero code change).
* **PPTX deck** — ``python-pptx``, fully local with zero system libraries
  (:func:`build_management_deck`; generalized to all 10 SPEC §11 slide types in a later task).

The SPEC §4.6 "Markdown → weasyprint" step is resolved to Jinja2→HTML→weasyprint because no
Markdown→HTML library is permitted by SPEC §6 / RULE 3 (documented in PROGRESS).
"""

from __future__ import annotations

import base64
import os
import re
import sys
from datetime import date
from functools import lru_cache
from pathlib import Path
from typing import Any

from jinja2 import Environment, FileSystemLoader, select_autoescape
from markupsafe import Markup, escape

from core import (
    blocks,
    charts_image,
    composer,
    deck_director,
    design,
    font_embed,
    imagery,
    layout,
    templates,
    visuals,
)
from core.artifact_framework import (
    brief_frame_context,
    deck_frame_label,
    framework_marker,
    prepare_brief_for_render,
    prepare_deck_for_render,
)
from core.config import AppConfig, ColorsConfig, StyleConfig
from core.deck_director import SlidePlan
from models.deck import Archetype, DeckStructure, SlideContent, SlideType
from models.diagram import DiagramNode, DiagramSpec, DiagramType
from models.research import ResearchReport
from models.synthesis import AnalysisOutput
from models.task import Audience, Language, TaskType
from models.visuals import ChartSpec


class ExportError(Exception):
    """Base class for output-rendering failures."""


class DeckBuildError(ExportError):
    """python-pptx is unavailable (fail fast with fix instructions)."""


class PdfBuildError(ExportError):
    """weasyprint (or its GTK system libraries) is unavailable (fail fast with fix instructions)."""


_PPTX_FIX = (
    "Fix: install python-pptx into the venv:\n  .venv\\Scripts\\python -m pip install python-pptx"
)
_PDF_FIX = (
    "PDF rendering needs WeasyPrint + the GTK/Pango runtime (PPTX + Excel work without it).\n"
    "Fix (Windows, one-time) — recommended via MSYS2 (auto-detected at C:\\msys64\\mingw64\\bin):\n"
    "  1. Install MSYS2 from https://www.msys2.org/ (keep the default C:\\msys64).\n"
    "  2. Open the 'MSYS2 MINGW64' terminal and run:\n"
    "       pacman -S mingw-w64-x86_64-pango\n"
    "     (pulls in cairo/glib/harfbuzz — everything WeasyPrint needs).\n"
    "  3. Reopen this terminal, then re-run.\n"
    "Alternative: install the GTK3 runtime (open the repo, then Releases):\n"
    "  https://github.com/tschoonj/GTK-for-Windows-Runtime-Installer  (tick 'set up PATH').\n"
    "If GTK lives elsewhere, set output.gtk_runtime_path in config.yaml to its 'bin' dir.\n"
    "Docs: https://www.gtk.org/docs/installations/windows/ · "
    "https://doc.courtbouillon.org/weasyprint/stable/first_steps.html#troubleshooting"
)


def _slug(text: str) -> str:
    """Make a short, filesystem-safe slug from a title."""
    cleaned = re.sub(r"[^a-z0-9]+", "_", text.lower()).strip("_")
    return (cleaned or "briefing")[:50]


def _sentences(text: str, cap: int = 6) -> list[str]:
    """Split a body into a few short bullet points (by line, then by sentence)."""
    raw = [seg.strip(" -•\t") for seg in re.split(r"\n+", text) if seg.strip()]
    bullets: list[str] = []
    for segment in raw:
        for part in re.split(r"(?<=[.!?])\s+", segment):
            part = part.strip()
            if part:
                bullets.append(part)
    return bullets[:cap] or (["—"] if not text.strip() else [text.strip()[:200]])


# ----------------------------------------------------------------- GTK runtime (WeasyPrint)
def _gtk_candidate_dirs(config: AppConfig) -> list[Path]:
    """Candidate GTK-runtime ``bin`` dirs (config > env > standard Windows locations)."""
    candidates: list[str] = []
    configured = config.output.gtk_runtime_path
    if configured:
        candidates.append(configured)
    env_dir = os.environ.get("WEASYPRINT_DLL_DIRECTORIES")
    if env_dir:
        candidates.extend(env_dir.split(os.pathsep))
    local = os.environ.get("LOCALAPPDATA", "")
    candidates += [
        r"C:\Program Files\GTK3-Runtime Win64\bin",
        os.path.join(local, "Programs", "GTK3-Runtime Win64", "bin") if local else "",
        r"C:\msys64\mingw64\bin",
    ]
    return [Path(c) for c in candidates if c]


def _ensure_gtk_dll_dir(config: AppConfig) -> None:
    """Put a real GTK runtime's ``bin`` first on the Windows DLL search path (idempotent).

    WeasyPrint needs Pango/Cairo/GObject. On Windows the Tesseract folder may ship an
    incompatible ``libgobject`` earlier on PATH, so a found GTK runtime is forced ahead of it
    (both via ``PATH`` and :func:`os.add_dll_directory`). No-op off Windows / when none is found.
    """
    if sys.platform != "win32":
        return
    for path in _gtk_candidate_dirs(config):
        if path.is_dir() and (path / "libgobject-2.0-0.dll").is_file():
            bin_str = str(path)
            current = os.environ.get("PATH", "")
            if not current.lower().startswith(bin_str.lower()):
                os.environ["PATH"] = bin_str + os.pathsep + current
            add_dll = getattr(os, "add_dll_directory", None)
            if callable(add_dll):
                try:
                    add_dll(bin_str)
                except OSError:  # pragma: no cover - defensive
                    pass
            return


# ----------------------------------------------------------------- brief (Jinja2 → HTML → PDF)
# Per-template labels (bottom-line label · type label), bilingual as (DE, EN).
_BRIEF_META: dict[str, dict[str, tuple[str, str]]] = {
    "competitor_brief.md.j2": {
        "bl": ("Kernaussage", "Executive Summary"),
        "type": ("Wettbewerbsanalyse", "Competitive Intelligence"),
    },
    "decision_brief.md.j2": {
        "bl": ("Empfehlung", "Recommendation"),
        "type": ("Entscheidungsanalyse", "Decision Analysis"),
    },
    "market_overview.md.j2": {
        "bl": ("Kernaussage", "Bottom Line"),
        "type": ("Marktanalyse", "Market Intelligence"),
    },
    "board_update.md.j2": {
        "bl": ("Kernaussage", "Bottom Line"),
        "type": ("Management-Briefing", "Management Brief"),
    },
    "document_synthesis.md.j2": {
        "bl": ("Kernaussage", "Bottom Line"),
        "type": ("Dokumentenanalyse", "Document Analysis"),
    },
    "adhoc_brief.md.j2": {
        "bl": ("Kernaussage", "Bottom Line Up Front"),
        "type": ("Kurzanalyse", "Quick Intel"),
    },
}

_DEFAULT_BRIEF_TEMPLATE = "adhoc_brief.md.j2"

# Task type → brief template (SPEC §10 T-1..T-6).
_BRIEF_TEMPLATES: dict[TaskType, str] = {
    TaskType.COMPETITOR_ANALYSIS: "competitor_brief.md.j2",
    TaskType.TARGET_SCREENING: "decision_brief.md.j2",
    TaskType.PARTNERSHIP_EVALUATION: "decision_brief.md.j2",
    TaskType.OPTION_COMPARISON: "decision_brief.md.j2",
    TaskType.STRATEGIC_INITIATIVE: "decision_brief.md.j2",
    TaskType.BUSINESS_CASE: "decision_brief.md.j2",
    TaskType.MARKET_RESEARCH: "market_overview.md.j2",
    TaskType.MARKET_ANALYSIS: "market_overview.md.j2",
    TaskType.FINANCIAL_BENCHMARK: "market_overview.md.j2",
    TaskType.BOARD_PREP: "board_update.md.j2",
    TaskType.MEETING_BRIEFING: "board_update.md.j2",
    TaskType.DOCUMENT_SYNTHESIS: "document_synthesis.md.j2",
    TaskType.INDUSTRY_NEWS: "adhoc_brief.md.j2",
    TaskType.ADHOC: "adhoc_brief.md.j2",
}

_TEMPLATES_DIR = Path(__file__).resolve().parent.parent / "templates" / "briefs"
_BULLET_RE = re.compile(r"^\s*[-*•]\s+")


def brief_template_for(task_type: TaskType) -> str:
    """Return the brief template filename for a task type (SPEC §10)."""
    return _BRIEF_TEMPLATES.get(task_type, _DEFAULT_BRIEF_TEMPLATE)


@lru_cache(maxsize=1)
def _brief_env() -> Environment:
    """Build (and cache) the Jinja2 environment for the brief templates."""
    return Environment(
        loader=FileSystemLoader(str(_TEMPLATES_DIR)),
        autoescape=select_autoescape(default=True, default_for_string=True),
        trim_blocks=True,
        lstrip_blocks=True,
    )


def _body_to_html(text: str) -> Markup:
    """Convert a section body (plain text + optional bullet lines) into safe HTML.

    Blank lines separate blocks; lines starting with ``-``/``*``/``•`` become a bullet list,
    other lines a paragraph (single newlines → ``<br>``). Everything is HTML-escaped first.
    """
    blocks = re.split(r"\n\s*\n", text.strip())
    parts: list[str] = []
    for block in blocks:
        lines = [ln for ln in block.splitlines() if ln.strip()]
        if not lines:
            continue
        if all(_BULLET_RE.match(ln) for ln in lines):
            items = "".join(f"<li>{escape(_BULLET_RE.sub('', ln).strip())}</li>" for ln in lines)
            parts.append(f"<ul>{items}</ul>")
        else:
            joined = Markup("<br>").join(escape(ln.strip()) for ln in lines)
            parts.append(f"<p>{joined}</p>")
    return Markup("".join(parts) or "<p>—</p>")


def _logo_data_uri(config: AppConfig) -> str | None:
    """Return the Neura logo as a base64 PNG data URI for HTML embedding (or ``None``)."""
    logo = Path(config.output.logo_path)
    if not (config.output.include_logo and logo.is_file()):
        return None
    data = base64.b64encode(logo.read_bytes()).decode("ascii")
    return f"data:image/png;base64,{data}"


# ----------------------------------------------------------------- editorial design tokens (PDF)
_FONT_EXTS = (".ttf", ".otf", ".woff2", ".woff")
_FONT_FORMATS = {".ttf": "truetype", ".otf": "opentype", ".woff2": "woff2", ".woff": "woff"}


def _font_face_css(style: StyleConfig) -> str:
    """``@font-face`` blocks for any shipped OFL TTFs found in ``style.fonts_dir`` (else empty).

    The CSS font stacks always keep their system-font fallbacks (``core.design.*_stack``), so a
    missing ``fonts_dir`` (the installer ships in Block 4) degrades to system fonts rather than
    breaking the brief (REQ-1/2). A face is emitted only for a file actually present, matched to a
    family by its normalized name appearing in the filename — robust to the installer's naming.
    """
    fonts_dir = Path(style.fonts_dir)
    if not fonts_dir.is_dir():
        return ""
    # Embed both the legacy configured families and every family of the active type-theme (so the
    # serif/expressive faces a deck/brief actually uses are subset-embedded into the PDF).
    theme_families = tuple(design.deck_fonts(style).values())
    families = dict.fromkeys(
        f
        for f in (
            style.serif_font,
            style.grotesk_font,
            style.body_font,
            style.mono_font,
            *theme_families,
        )
        if f
    )
    candidates = sorted(p for p in fonts_dir.iterdir() if p.suffix.lower() in _FONT_EXTS)
    blocks: list[str] = []
    for family in families:
        norm = re.sub(r"[\s_-]+", "", family).lower()
        match = next(
            (p for p in candidates if norm in re.sub(r"[\s_-]+", "", p.stem).lower()), None
        )
        if match is None:
            continue
        url = match.resolve().as_uri()
        fmt = _FONT_FORMATS[match.suffix.lower()]
        # Minimal, widely-supported descriptors only: a `font-weight` *range* (CSS Fonts 4) is
        # rejected by WeasyPrint, so it is omitted — the variable font still embeds and heavier
        # requests (bold headlines) render from the same face (synthesized when needed).
        blocks.append(
            f'@font-face {{ font-family: "{family}"; src: url("{url}") format("{fmt}"); }}'
        )
    return "\n".join(blocks)


def _cover_band_svg(colors: ColorsConfig) -> Markup:
    """A slim full-width luminous warm→cool gradient band for the editorial brief header (DNA 5)."""
    gradient = design.linear_gradient_svg(
        "porter-coverband", design.depth_gradient_stops(colors), x1=0.0, y1=0.0, x2=1.0, y2=0.0
    )
    return Markup(  # gradient is built from config hex values, not user input
        '<svg class="cover-band" viewBox="0 0 100 6" preserveAspectRatio="none" '
        'xmlns="http://www.w3.org/2000/svg" aria-hidden="true">'
        f"<defs>{gradient}</defs>"
        '<rect x="0" y="0" width="100" height="6" fill="url(#porter-coverband)"/></svg>'
    )


def _brief_chart_svg(spec: ChartSpec, colors: ColorsConfig, style: StyleConfig) -> Markup | None:
    """Render a section's chart spec to an inline ``<svg>`` for the brief (fail-open → ``None``).

    Uses the same hand-built SVG engine as the locked render tech (no raster, no new dep). The spec
    is renderable by construction (``ChartSpec`` validates structure); any rendering quirk degrades
    to no chart rather than breaking the brief (REQ-5).
    """
    try:
        svg = visuals.render_chart_svg(spec, colors, style, on_dark=False)
    except Exception:  # noqa: BLE001 — charts are decorative; never break a brief
        return None
    return Markup(svg)  # themed SVG built from the validated spec, not user markup


def _brief_context(
    analysis: AnalysisOutput,
    config: AppConfig,
    template_name: str,
    *,
    task_type: TaskType,
    audience: Audience | None = None,
    research_report: ResearchReport | None = None,
) -> dict[str, Any]:
    """Assemble the Jinja context for a brief from the analysis (bilingual labels).

    Besides the content (title/bottom line/sections/sources) this injects the Porter Editorial
    design tokens consumed by the templates: the multi-font CSS stacks (serif/grotesk/body/mono),
    the two-tone headline split (DNA 7), source-grounded telemetry chips (DNA 6), the editorial
    intensity flag + luminous cover band (DNA 5), the chart palette, and per-section inline chart
    SVGs (the locked PDF render tech). Renderers never call the LLM and invent no data.
    """
    is_de = analysis.language == Language.DE
    idx = 0 if is_de else 1
    colors = config.output.colors
    style = config.output.style
    editorial = design.is_editorial(style)
    meta = _BRIEF_META.get(template_name, _BRIEF_META[_DEFAULT_BRIEF_TEMPLATE])
    meta_line = (
        f"{date.today().isoformat()} · {meta['type'][idx]} · {analysis.language.value.upper()}"
    )
    sections = _brief_sections(analysis, colors, style)
    title_before, title_hl, title_after = design.split_for_highlight(analysis.title)
    sources = [
        {
            "url": s.url,
            "title": s.title or "",
            "date": s.date or "",
            "tier": s.tier.value if s.tier else "",
        }
        for s in analysis.sources
    ]
    context = {
        "c": colors,
        "title": analysis.title,
        "title_before": title_before,
        "title_hl": title_hl,
        "title_after": title_after,
        "meta": meta_line,
        "bl_label": meta["bl"][idx],
        "bottom_line": analysis.bottom_line,
        "sections": sections,
        "sources_label": "Quellen" if is_de else "Sources",
        "sources": sources,
        "note": "",
        "footer_note": "NEURA · " + ("Management-Briefing" if is_de else "Management Brief"),
        "logo_data_uri": _logo_data_uri(config),
        # --- editorial design tokens ---
        # Font stacks are CSS values (quoted family names) injected into <style>, so they must NOT
        # be HTML-escaped — Markup keeps the quotes intact (config-derived RULE 4; never user text).
        "editorial": editorial,
        "font_serif": Markup(design.serif_stack(style)),
        "font_grotesk": Markup(design.grotesk_stack(style)),
        "font_body": Markup(design.body_stack(style)),
        "font_mono": Markup(design.mono_stack(style)),
        "font_faces_css": Markup(_font_face_css(style)),
        "palette": design.chart_series_colors(colors),
        "spot": colors.coral,  # two-tone highlight: coral on the cream canvas (DNA 7)
        "cover_band_svg": _cover_band_svg(colors) if editorial else None,
        "telemetry": design.telemetry_chips(research_report, analysis.language),
    }
    context.update(brief_frame_context(analysis, task_type=task_type, audience=audience))
    return context


def _brief_sections(
    analysis: AnalysisOutput, colors: ColorsConfig, style: StyleConfig
) -> list[dict[str, Any]]:
    """Build the section render contexts, embedding an inline chart SVG for each ``.visual``.

    A section carrying a :class:`~models.visuals.ChartSpec` gets a themed inline ``<svg>`` until the
    ``style.max_charts_per_brief`` budget is spent (and only while ``charts_enabled``); the rest
    render text-only. Deterministic + fail-open — Block 4 owns where briefs source their specs.
    """
    charts_used = 0
    out: list[dict[str, Any]] = []
    for section in analysis.sections:
        chart_svg: Markup | None = None
        if (
            section.visual is not None
            and style.charts_enabled
            and charts_used < style.max_charts_per_brief
        ):
            chart_svg = _brief_chart_svg(section.visual, colors, style)
            if chart_svg is not None:
                charts_used += 1
        out.append(
            {
                "heading": section.heading,
                "body_html": _body_to_html(section.body),
                "chart_svg": chart_svg,
            }
        )
    return out


def render_brief_html(
    analysis: AnalysisOutput,
    config: AppConfig,
    *,
    task_type: TaskType = TaskType.DOCUMENT_SYNTHESIS,
    audience: Audience | None = None,
    research_report: ResearchReport | None = None,
) -> str:
    """Render the analysis into a Porter Editorial HTML brief (pure — no WeasyPrint, testable).

    When ``research_report`` is supplied, source-grounded telemetry chips (DNA 6) render in the
    header — the pipeline threads it in (Block 4). Otherwise no chips appear (never invented).
    """
    analysis = prepare_brief_for_render(analysis)
    template_name = brief_template_for(task_type)
    template = _brief_env().get_template(template_name)
    return template.render(
        **_brief_context(
            analysis,
            config,
            template_name,
            task_type=task_type,
            audience=audience,
            research_report=research_report,
        )
    )


def build_brief_pdf(
    analysis: AnalysisOutput,
    config: AppConfig,
    output_dir: str | Path,
    *,
    task_type: TaskType = TaskType.DOCUMENT_SYNTHESIS,
    audience: Audience | None = None,
    research_report: ResearchReport | None = None,
) -> Path:
    """Render the analysis into a Porter Editorial PDF brief (Jinja2 → HTML → WeasyPrint).

    The brief template is chosen from the task type (SPEC §10); ``research_report`` lights up the
    telemetry chips (mirrors :func:`build_deck`). Raises :class:`PdfBuildError` with exact install
    instructions if WeasyPrint / its GTK runtime is unavailable (the renderer is correct and works
    the moment GTK is present — no code change).
    """
    _ensure_gtk_dll_dir(config)
    try:
        from weasyprint import HTML
    except (ImportError, OSError) as exc:
        raise PdfBuildError(_PDF_FIX) from exc

    html = render_brief_html(
        analysis, config, task_type=task_type, audience=audience, research_report=research_report
    )
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    path = out / f"{date.today().isoformat()}_{_slug(analysis.title)}_brief.pdf"
    HTML(string=html).write_pdf(str(path))
    return path


# ----------------------------------------------------------------- PPTX (python-pptx)
_SLIDE_W_IN = 13.333  # 16:9 widescreen
_SLIDE_H_IN = 7.5
# SWOT quadrant fills (Strengths/Opportunities = positive, Weaknesses/Threats = risk).
_SWOT_FILLS = ("excel_positive", "excel_negative", "excel_positive", "excel_negative")
# A *hero* metric must carry a currency or a unit — a bare integer like the "1" in "Focus Area 1" is
# a list ordinal, NOT a metric, so it never gets promoted to a giant accent numeral (the v3 bug).
# Units are ordered longest-first with a trailing word boundary (so "months" is not sliced to "m").
_METRIC_RE = re.compile(
    r"(?:(?:EUR|USD|GBP)\s?|[$€£]\s?)\d[\d.,]*"  # currency-prefixed amount
    r"|"
    # number + unit (units longest-first, alpha ones need a \b so "months" is not sliced to "m")
    r"\d[\d.,]*\s?(?:%|x|(?:million|billion|months?|years?|mio|bn|EUR|USD|GBP|k|m|b)\b)",
    re.IGNORECASE,
)


class _DeckRenderer:
    """Render a :class:`~models.deck.DeckStructure` into a Porter Editorial python-pptx deck.

    Holds the shared pptx primitives (colors, multi-font text, bullets, tables, rounded boxes,
    the bottom-right logo) and dispatches per :class:`~models.deck.SlideType`, so all 10 SPEC §11
    slide types share one implementation (no per-type duplication, SPEC §11 + output_playbook).

    The Editorial system (``core/design.py``) drives the look: a cream ``paper`` canvas for content
    slides, a dramatic ``canvas_dark`` cover/divider with a luminous gradient + focal glow (only on
    editorial intensity), a grotesk/body/mono multi-font system, two-tone headlines, DNA-4 system
    cards, source-grounded telemetry chips, and native editable charts for a slide's ``.visual``.
    Renderers never call the LLM. ``restrained`` intensity drops the gradient/glow (board-safe).
    """

    def __init__(
        self,
        config: AppConfig,
        language: Language,
        research_report: ResearchReport | None = None,
    ) -> None:
        """Build an empty 16:9 presentation and cache the pptx primitives + Editorial tokens."""
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
            from pptx.oxml import parse_xml
            from pptx.oxml.ns import nsdecls, qn
            from pptx.util import Inches, Pt
        except ImportError as exc:  # pragma: no cover - python-pptx is a declared dependency
            raise DeckBuildError(f"python-pptx is not installed.\n{_PPTX_FIX}") from exc

        self._RGBColor = RGBColor
        self._MSO_SHAPE = MSO_SHAPE
        self._PP_ALIGN = PP_ALIGN
        self._MSO_ANCHOR = MSO_ANCHOR
        self._MSO_AUTO_SIZE = MSO_AUTO_SIZE
        self._Inches = Inches
        self._Pt = Pt
        self._parse_xml = parse_xml
        self._nsdecls = nsdecls
        self._qn = qn
        self.colors = config.output.colors
        self.style = config.output.style
        self.language = language
        self.research_report = research_report
        self._editorial = design.is_editorial(self.style)
        self.editorial = self._editorial  # public alias for the Surface protocol (core/blocks)
        fonts = design.deck_fonts(self.style)
        self.fonts = fonts  # role → family, consumed by core/blocks via the Surface protocol
        self._font_display = fonts["display"]
        self._font_body = fonts["body"]
        self._font_mono = fonts["mono"]
        self._font_serif = fonts.get("serif", fonts["display"])
        self._font_statement = fonts.get("statement", fonts["display"])
        self._charts_used = 0
        self._max_charts = self.style.max_charts_per_deck
        self._diagrams_used = 0
        self._max_diagrams = self.style.max_diagrams_per_deck
        self._slide_no = 0
        # Per-slide canvas state (set by ``_frame`` / ``_title_slide``): foreground color and
        # whether the slide sits on the dark editorial canvas (so headline/labels stay legible).
        self._slide_fg = self.colors.ink
        self._slide_on_dark = False
        self.prs = Presentation()
        self.prs.slide_width = Inches(_SLIDE_W_IN)
        self.prs.slide_height = Inches(_SLIDE_H_IN)
        self._blank = self.prs.slide_layouts[6]
        logo = Path(config.output.logo_path)
        self._show_logo = config.output.include_logo and logo.is_file()
        self._logo = str(logo)
        light = config.output.logo_path_light
        self._logo_light = str(light) if light and Path(light).is_file() else None
        self._imagery_dir = config.output.imagery_dir

    def rgb(self, hex_color: str) -> Any:
        """A pptx ``RGBColor`` from a ``#rrggbb`` string (config-driven palette)."""
        return self._RGBColor.from_string(hex_color.lstrip("#").upper())  # type: ignore[no-untyped-call]

    def _add_logo(self, slide: Any) -> None:
        """Place the Neura logo bottom-right: width-capped, inset with a margin, light on dark.

        The width cap stops a wide logo from sprawling across the slide, and the bottom-right inset
        keeps it clear of the bottom-left page number (the v3 output stacked them). On a dark canvas
        a configured light logo variant is used so the mark does not vanish (Block 5.0).
        """
        if not self._show_logo:
            return
        inches = self._Inches
        path = self._logo_light if (self._slide_on_dark and self._logo_light) else self._logo
        pic = slide.shapes.add_picture(path, 0, 0, height=inches(0.32))
        max_w = inches(1.5)
        if pic.width > max_w:
            ratio = max_w / pic.width
            pic.width = int(pic.width * ratio)
            pic.height = int(pic.height * ratio)
        pic.left = self.prs.slide_width - pic.width - inches(0.3)
        pic.top = inches(_SLIDE_H_IN) - pic.height - inches(0.24)

    def _text(
        self,
        slide: Any,
        text: str,
        left: Any,
        top: Any,
        width: Any,
        height: Any,
        size: int,
        color: str,
        *,
        bold: bool = False,
        align: Any = None,
        anchor: Any = None,
        font: str | None = None,
    ) -> Any:
        """Add a single-paragraph text box and return it (multi-font: grotesk/body/mono).

        ``font`` forces a font role (e.g. mono micro-labels); when ``None`` the size picks grotesk
        display for headings (>=24pt) and the body font otherwise.
        """
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        if anchor is not None:
            frame.vertical_anchor = anchor
        para = frame.paragraphs[0]
        para.text = text
        para.alignment = align if align is not None else self._PP_ALIGN.LEFT
        run = para.runs[0] if para.runs else para.add_run()
        run.font.size = self._Pt(size)
        run.font.bold = bold
        run.font.name = font or (self._font_display if size >= 24 else self._font_body)
        run.font.color.rgb = self.rgb(color)
        return box

    def _accent(self, slide_type: SlideType) -> str:
        """Semantic accent color for the slide frame."""
        if slide_type == SlideType.RECOMMENDATION:
            return self.colors.artifact_teal
        if slide_type in {SlideType.SWOT, SlideType.FINANCIAL_OVERVIEW}:
            return self.colors.artifact_gold
        if slide_type == SlideType.APPENDIX:
            return self.colors.charcoal
        return self.colors.accent_cyan

    def _rect(
        self,
        slide: Any,
        left: Any,
        top: Any,
        width: Any,
        height: Any,
        fill: str,
        *,
        line: str | None = None,
    ) -> Any:
        """Add a plain rectangle with optional line color."""
        shape = slide.shapes.add_shape(self._MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.rgb(fill)
        if line:
            shape.line.color.rgb = self.rgb(line)
            shape.line.width = self._Pt(0.75)
        else:
            shape.line.fill.background()
        return shape

    # --- editorial depth (luminous dark canvas) ------------------------------------------
    def _apply_gradient(self, shape: Any, stops: list[tuple[float, str]]) -> None:
        """Replace a shape's fill with a multi-stop diagonal gradient (warm→cool depth, DNA 5).

        Built as raw DrawingML so all of ``depth_gradient_stops`` is honored (python-pptx's high
        level API only exposes two stops). The ``<a:gradFill>`` element is parsed first, so if the
        XML build fails the shape keeps its original solid fill (fail-open).
        """
        gs = "".join(
            f'<a:gs pos="{int(round(offset * 100000))}">'
            f'<a:srgbClr val="{hex_color.lstrip("#").upper()}"/></a:gs>'
            for offset, hex_color in stops
        )
        xml = (
            f"<a:gradFill {self._nsdecls('a')}><a:gsLst>{gs}</a:gsLst>"
            f'<a:lin ang="2700000" scaled="1"/></a:gradFill>'
        )
        grad = self._parse_xml(xml)
        spset = shape.fill._xPr
        fill_kinds = ("solidFill", "gradFill", "noFill", "blipFill", "pattFill", "grpFill")
        fill_tags = {self._qn(f"a:{kind}") for kind in fill_kinds}
        for child in list(spset):
            if child.tag in fill_tags:
                spset.remove(child)
        line_el = spset.find(self._qn("a:ln"))
        if line_el is not None:
            line_el.addprevious(grad)
        else:
            spset.append(grad)

    def _set_alpha(self, shape: Any, opacity_pct: int) -> None:
        """Set a solid-filled shape's opacity (0..100) via the DrawingML ``<a:alpha>`` element."""
        spset = shape.fill._xPr
        solid = spset.find(self._qn("a:solidFill"))
        if solid is None:
            return
        srgb = solid.find(self._qn("a:srgbClr"))
        if srgb is None:
            return
        for existing in srgb.findall(self._qn("a:alpha")):
            srgb.remove(existing)
        alpha = f'<a:alpha {self._nsdecls("a")} val="{int(opacity_pct * 1000)}"/>'
        srgb.append(self._parse_xml(alpha))

    def _soft_edge(self, shape: Any, radius_in: float) -> None:
        """Blur a shape's edge (turns a flat oval into a soft luminous glow).

        ``CT_ShapeProperties`` allows only ONE ``<a:effectLst>`` — and ``shadow.inherit = False``
        already creates an empty one — so the soft edge is added *into* the existing list (a second
        ``effectLst`` is a schema violation PowerPoint silently repairs away).
        """
        radius = int(radius_in * 914400)  # inches → EMU
        spPr = shape._element.spPr
        effect_lst = spPr.find(self._qn("a:effectLst"))
        if effect_lst is None:
            effect_lst = self._parse_xml(f"<a:effectLst {self._nsdecls('a')}/>")
            line_el = spPr.find(self._qn("a:ln"))
            if line_el is not None:
                line_el.addnext(effect_lst)
            else:
                spPr.append(effect_lst)
        effect_lst.append(self._parse_xml(f'<a:softEdge {self._nsdecls("a")} rad="{radius}"/>'))

    def _add_glow(self, slide: Any) -> None:
        """One warm, soft focal glow upper-right of a dark canvas (DNA 5: a luminous accent)."""
        inches = self._Inches
        oval = slide.shapes.add_shape(
            self._MSO_SHAPE.OVAL, inches(8.9), inches(-1.7), inches(6.2), inches(6.2)
        )
        oval.line.fill.background()
        oval.shadow.inherit = False
        oval.fill.solid()
        oval.fill.fore_color.rgb = self.rgb(design.glow_color(self.colors))
        self._set_alpha(oval, 24)
        self._soft_edge(oval, 1.1)

    def _paint_dark_canvas(self, slide: Any, *, glow: bool = True) -> None:
        """Paint a dramatic dark editorial canvas (cover/divider moment).

        Always lays down a solid ``canvas_dark`` background; on editorial intensity it adds a
        full-bleed luminous warm→cool gradient and one warm focal glow on top. All depth is
        decorative and wrapped fail-open, so a deck never breaks because of a fill quirk.
        """
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = self.rgb(self.colors.canvas_dark)
        if not self._editorial:
            return
        try:
            inches = self._Inches
            wash = self._rect(
                slide,
                inches(0),
                inches(0),
                inches(_SLIDE_W_IN),
                inches(_SLIDE_H_IN),
                self.colors.canvas_dark,
            )
            wash.shadow.inherit = False
            self._apply_gradient(wash, design.depth_gradient_stops(self.colors))
            if glow:
                self._add_glow(slide)
        except Exception:  # noqa: BLE001 — depth is decorative; never break a deck (REQ-5)
            pass

    def _paint_cover_image(self, slide: Any, image_path: Any) -> None:
        """Lay a full-bleed brand image under a dark scrim so the knockout title reads (fail-open).

        On any image error the cover degrades to the luminous gradient canvas (REQ-5).
        """
        inches = self._Inches
        try:
            slide_w = inches(_SLIDE_W_IN)
            slide_h = inches(_SLIDE_H_IN)
            # Cover-fit: scale to fill the 16:9 frame preserving aspect ratio (crop the overflow),
            # centered — so a square brand photo is never squished. Add native, then resize/center.
            pic = slide.shapes.add_picture(str(image_path), 0, 0)
            if pic.width and pic.height:
                scale = max(slide_w / pic.width, slide_h / pic.height)
                pic.width = int(pic.width * scale)
                pic.height = int(pic.height * scale)
            pic.left = int((slide_w - pic.width) / 2)
            pic.top = int((slide_h - pic.height) / 2)
            scrim = self._rect(
                slide, inches(0), inches(0), slide_w, slide_h, self.colors.canvas_dark
            )
            scrim.shadow.inherit = False
            self._set_alpha(scrim, 55)
        except Exception:  # noqa: BLE001 — a bad image never loses the cover (REQ-5)
            self._paint_dark_canvas(slide, glow=True)

    def _frame(self, slide: Any, sc: SlideContent, *, canvas_hex: str | None = None) -> None:
        """Apply the mandatory Porter Editorial frame to a non-title slide on a given canvas.

        ``canvas_hex`` is the design-director's chosen background (cream / sand / cream-hi, or the
        dark divider for the recommendation). When omitted, content sits on cream and the
        RECOMMENDATION slide becomes the dramatic dark divider. Sets the per-slide foreground color
        so headlines/labels stay legible on whichever canvas is used.
        """
        self._slide_no += 1
        if canvas_hex is None:
            canvas = (
                self.colors.canvas_dark
                if sc.slide_type == SlideType.RECOMMENDATION
                else self.colors.paper
            )
        else:
            canvas = canvas_hex
        self._paint_solid_or_dark(slide, canvas)
        self._frame_chrome(slide, sc, number=self._slide_no)

    def _paint_solid_or_dark(self, slide: Any, canvas: str) -> None:
        """Paint a content canvas (solid light, or the luminous dark) + set the per-slide fg."""
        on_dark = design.luminance(canvas) < 0.55
        self._slide_on_dark = on_dark
        self._slide_fg = design.contrast_text(canvas, self.colors)
        if on_dark:
            self._paint_dark_canvas(slide, glow=True)
        else:
            background = slide.background
            background.fill.solid()
            background.fill.fore_color.rgb = self.rgb(canvas)

    def _frame_chrome(self, slide: Any, sc: SlideContent, *, number: int) -> None:
        """Draw the editorial chrome (spine · footer rule · label · page no · telemetry · logo).

        Canvas-agnostic: it reads ``self._slide_on_dark`` (set by whatever painted the background),
        so both the legacy ``_frame`` and the Block-2 composition path share one chrome.
        """
        inches = self._Inches
        muted = self.colors.light_surface if self._slide_on_dark else self.colors.charcoal
        accent = self._accent(sc.slide_type)
        self._rect(
            slide, inches(0), inches(0), inches(0.16), inches(_SLIDE_H_IN), self.colors.canvas_dark
        )
        self._rect(slide, inches(0.16), inches(0), inches(0.045), inches(_SLIDE_H_IN), accent)
        self._rect(slide, inches(0.65), inches(6.93), inches(10.9), inches(0.012), muted)
        label = f"{deck_frame_label(self.language, sc.slide_type).upper()}  ·  {framework_marker()}"
        self._text(
            slide,
            label,
            inches(0.65),
            inches(0.16),
            inches(8.8),
            inches(0.25),
            8,
            muted,
            bold=True,
            font=self._font_mono,
        )
        # Page number bottom-LEFT (the bottom-right corner is reserved for the logo, so they never
        # collide — the v3 output stacked the logo on top of the page number).
        self._text(
            slide,
            f"{number:02d}",
            inches(0.65),
            inches(6.98),
            inches(0.5),
            inches(0.2),
            8,
            muted,
            bold=True,
            align=self._PP_ALIGN.LEFT,
            font=self._font_mono,
        )
        if self.style.telemetry_chips:
            self._telemetry_footer(slide, color=muted)
        self._add_logo(slide)

    def _telemetry_footer(self, slide: Any, *, color: str) -> None:
        """Render source-grounded telemetry as mono outline pills on the bottom rail (DNA 6).

        Chips come from :func:`core.design.telemetry_chips` (only real counts/dates from the
        research report); nothing renders when no report was threaded in. Never invents data.
        """
        chips = design.telemetry_chips(self.research_report, self.language)
        if not chips:
            return
        inches = self._Inches
        x = 0.65
        for chip in chips:
            width = 0.30 + 0.072 * len(chip)
            if x + width > 11.35:  # keep clear of the page number bottom-right
                break
            self._rounded(
                slide,
                inches(x),
                inches(6.99),
                inches(width),
                inches(0.28),
                self.colors.paper,
                line=color,
            )
            self._text(
                slide,
                chip,
                inches(x),
                inches(7.0),
                inches(width),
                inches(0.26),
                8,
                color,
                bold=True,
                align=self._PP_ALIGN.CENTER,
                anchor=self._MSO_ANCHOR.MIDDLE,
                font=self._font_mono,
            )
            x += width + 0.12

    def _metric_token(self, text: str) -> str | None:
        """Extract a visible number token from a line, if one exists."""
        match = _METRIC_RE.search(text)
        if not match:
            return None
        return " ".join(match.group(0).split())[:14]

    def _short(self, text: str, limit: int = 165) -> str:
        """Keep card text in fixed boxes, truncating at a sentence/word boundary (never mid-word).

        Strips any leaked inline Markdown, then prefers to end on a completed sentence inside the
        budget; failing that it cuts at the last whole word and adds an ellipsis — so a clipped
        bullet never reads as a dangling fragment like "… the most" (the v4 complaint).
        """
        cleaned = design.strip_inline_markdown(text)
        if len(cleaned) <= limit:
            return cleaned
        window = cleaned[:limit]
        sentence_end = max(window.rfind(". "), window.rfind("! "), window.rfind("? "))
        if sentence_end >= limit * 0.6:
            return window[: sentence_end + 1].rstrip()
        word_cut = window.rsplit(" ", 1)[0].rstrip(" ,;:—–-")
        return (word_cut or window).rstrip() + "…"

    def _body_callout(
        self, slide: Any, text: str, *, top: float = 1.55, fill: str | None = None
    ) -> None:
        """Render a high-contrast one-message callout."""
        inches = self._Inches
        box = self._rounded(
            slide,
            inches(0.72),
            inches(top),
            inches(11.75),
            inches(1.18),
            fill or self.colors.canvas_dark,
        )
        frame = box.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = self._MSO_ANCHOR.MIDDLE
        frame.margin_left = inches(0.22)
        frame.margin_right = inches(0.22)
        para = frame.paragraphs[0]
        para.text = self._short(text, 190)
        para.alignment = self._PP_ALIGN.CENTER
        run = para.runs[0] if para.runs else para.add_run()
        run.font.size = self._Pt(20)
        run.font.bold = True
        run.font.name = self._font_display
        run.font.color.rgb = self.rgb(self.colors.white)

    def _signal_cards(
        self, slide: Any, bullets: list[str], *, top: float = 1.75, max_cards: int = 4
    ) -> None:
        """Render support points as DNA-4 'system cards' (one accent per card, indexed, arrowed)."""
        palette = design.chart_series_colors(self.colors)
        lines = [self._short(line) for line in (bullets or ["-"]) if str(line).strip()]
        lines = lines[:max_cards] or ["-"]
        # With four long bullets the cards get too short to breathe — drop to three for legibility.
        if len(lines) == 4 and max(len(s) for s in lines) > 150:
            lines = lines[:3]
        count = len(lines)
        cols = 1 if count == 1 else 2
        rows = (count + cols - 1) // cols
        gap = 0.28
        card_w = 11.75 if cols == 1 else 5.73
        available_h = max(1.1, 6.65 - top)
        card_h = min(2.15, (available_h - gap * (rows - 1)) / rows)

        for idx, line in enumerate(lines):
            row, col = divmod(idx, cols)
            left = 0.72 + col * (card_w + gap)
            y = top + row * (card_h + gap)
            self._system_card(
                slide,
                line,
                left=left,
                y=y,
                width=card_w,
                height=card_h,
                accent=palette[idx % len(palette)],
                index=idx + 1,
            )

    def _system_card(
        self,
        slide: Any,
        line: str,
        *,
        left: float,
        y: float,
        width: float,
        height: float,
        accent: str,
        index: int,
    ) -> None:
        """One Editorial system card: white block, accent spine, big metric (or index), body, tag.

        Cards are always white so they read on both the cream and the dark canvas; the accent is the
        single color block per card (DNA 4). A metric token is surfaced big top-left when the line
        carries a number (DNA 6); otherwise a mono index marker. A corner-arrow sits top-right (DNA
        8) and a mono category tag bottom-left (DNA 3).
        """
        inches = self._Inches
        self._rounded(
            slide, inches(left), inches(y), inches(width), inches(height), self.colors.white
        )
        self._rounded(slide, inches(left), inches(y), inches(0.1), inches(height), accent)
        self._text(
            slide,
            "↗",  # ↗ corner-arrow
            inches(left + width - 0.52),
            inches(y + 0.12),
            inches(0.4),
            inches(0.32),
            13,
            accent,
            bold=True,
            align=self._PP_ALIGN.RIGHT,
            font=self._font_body,
        )
        token = self._metric_token(line)
        if token:
            self._text(
                slide,
                token,
                inches(left + 0.3),
                inches(y + 0.16),
                inches(width - 0.95),
                inches(0.6),
                24,
                accent,
                bold=True,
            )
            body_top = y + 0.84
        else:
            self._text(
                slide,
                f"{index:02d}",
                inches(left + 0.3),
                inches(y + 0.18),
                inches(1.2),
                inches(0.4),
                14,
                accent,
                bold=True,
                font=self._font_mono,
            )
            body_top = y + 0.64
        # Dynamic size + autofit so a long bullet shrinks to fit instead of spilling past the card
        # (the v3 overflow). The mono [NN] tag is only drawn when the card is tall enough to hold it
        # without colliding with the body text.
        show_tag = height >= 1.35
        body_size = 12 if len(line) <= 130 else 11 if len(line) <= 185 else 10
        reserve = 0.36 if show_tag else 0.14
        body_box = self._text(
            slide,
            line,
            inches(left + 0.3),
            inches(body_top),
            inches(width - 0.6),
            inches(max(0.4, height - (body_top - y) - reserve)),
            body_size,
            self.colors.ink,
        )
        body_box.text_frame.auto_size = self._MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        if show_tag:
            self._text(
                slide,
                f"[{index:02d}]",
                inches(left + 0.3),
                inches(y + height - 0.30),
                inches(1.6),
                inches(0.22),
                8,
                self.colors.charcoal,
                bold=True,
                font=self._font_mono,
            )

    def _compact_list(self, slide: Any, items: list[str], *, top: float = 1.75) -> None:
        """Render appendix/source lines in two dense, readable columns (mono, numbered references).

        Items arrive pre-numbered ("07  domain — title"), so the mono face + the leading index read
        as one consistent reference list — every appendix page in the same style (no stray bullet).
        """
        inches = self._Inches
        lines = [self._short(item, 116) for item in (items or ["-"])][:18]
        for idx, line in enumerate(lines):
            col = idx % 2
            row = idx // 2
            left = 0.75 + col * 5.9
            y = top + row * 0.52
            self._text(
                slide,
                line,
                inches(left),
                inches(y),
                inches(5.55),
                inches(0.42),
                9,
                self.colors.charcoal,
                font=self._font_mono,
            )

    def _spot_color(self) -> str:
        """Two-tone highlight color for the current canvas (warm gold on dark, coral on cream)."""
        return self.colors.artifact_gold if self._slide_on_dark else self.colors.coral

    def _headline_two_tone(
        self,
        slide: Any,
        text: str,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        size: int,
        base_color: str,
        accent_color: str,
    ) -> Any:
        """Render a headline in one textbox with a single key token in an accent run (DNA 7).

        Uses :func:`core.design.split_for_highlight` to pick the token (number, else proper noun);
        the wording is never changed — only one span is recolored. Falls back to a plain run.
        """
        inches = self._Inches
        before, token, after = design.split_for_highlight(text)
        box = slide.shapes.add_textbox(inches(left), inches(top), inches(width), inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        para = frame.paragraphs[0]
        para.alignment = self._PP_ALIGN.LEFT

        def _run(segment: str, color: str) -> None:
            run = para.add_run()
            run.text = segment
            run.font.size = self._Pt(size)
            run.font.bold = True
            run.font.name = self._font_display
            run.font.color.rgb = self.rgb(color)

        if token:
            if before:
                _run(before, base_color)
            _run(token, accent_color)
            if after:
                _run(after, base_color)
        else:
            _run(text, base_color)
        return box

    def _headline(self, slide: Any, text: str, *, accent: str) -> None:
        """Render the claim headline: a semantic accent bar + a two-tone grotesk headline."""
        inches = self._Inches
        self._rect(slide, inches(0.68), inches(0.58), inches(0.08), inches(0.72), accent)
        self._headline_two_tone(
            slide,
            text,
            left=0.9,
            top=0.5,
            width=11.4,
            height=1.0,
            size=25,
            base_color=self._slide_fg,
            accent_color=self._spot_color(),
        )

    def _bullets(self, slide: Any, bullets: list[str], *, top: float = 1.8) -> None:
        """Render a vertical list of bullets on a content slide."""
        inches, pt = self._Inches, self._Pt
        box = slide.shapes.add_textbox(
            inches(0.7), inches(top), inches(11.9), inches(_SLIDE_H_IN - top - 0.6)
        )
        frame = box.text_frame
        frame.word_wrap = True
        for idx, line in enumerate(bullets or ["—"]):
            para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            para.text = f"• {line}"
            para.space_after = pt(9)
            run = para.runs[0]
            run.font.size = pt(16)
            run.font.name = self._font_body
            run.font.color.rgb = self.rgb(self._slide_fg)

    def _rounded(
        self,
        slide: Any,
        left: Any,
        top: Any,
        width: Any,
        height: Any,
        fill: str,
        *,
        line: str | None = None,
    ) -> Any:
        """Add a borderless rounded rectangle filled with ``fill`` (clean Neura aesthetic)."""
        shape = slide.shapes.add_shape(self._MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.rgb(fill)
        if line:
            shape.line.color.rgb = self.rgb(line)
            shape.line.width = self._Pt(0.75)
        else:
            shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    # --- Editorial v4.0 shared helpers (archetypes + depth) ------------------------------
    def _soft_shadow(self, shape: Any, *, blur_in: float = 0.07, dist_in: float = 0.05) -> None:
        """Add a soft outer shadow into the shape's existing effectLst (editorial only, fail-open).

        Mirrors :meth:`_soft_edge`'s single-effectLst discipline (a second ``<a:effectLst>`` is a
        schema violation PowerPoint repairs away). Gives cards/tiles depth; never breaks a deck.
        """
        if not self._editorial:
            return
        try:
            blur = int(blur_in * 914400)
            dist = int(dist_in * 914400)
            spPr = shape._element.spPr
            effect_lst = spPr.find(self._qn("a:effectLst"))
            if effect_lst is None:
                effect_lst = self._parse_xml(f"<a:effectLst {self._nsdecls('a')}/>")
                line_el = spPr.find(self._qn("a:ln"))
                if line_el is not None:
                    line_el.addnext(effect_lst)
                else:
                    spPr.append(effect_lst)
            ink = self.colors.ink.lstrip("#").upper()
            shadow = (
                f'<a:outerShdw {self._nsdecls("a")} blurRad="{blur}" dist="{dist}" '
                f'dir="5400000" rotWithShape="0"><a:srgbClr val="{ink}">'
                f'<a:alpha val="28000"/></a:srgbClr></a:outerShdw>'
            )
            effect_lst.insert(0, self._parse_xml(shadow))
        except Exception:  # noqa: BLE001 — depth is decorative; never break a deck (REQ-5)
            pass

    def _paint_field(self, slide: Any, field_hex: str) -> None:
        """Paint a full-bleed saturated statement field; editorial adds a subtle duotone.

        Sets the per-slide foreground to a legible knockout on the field. The duotone (field → a
        slightly darker field) gives depth; restrained keeps a flat solid field.
        """
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = self.rgb(field_hex)
        self._slide_on_dark = design.luminance(field_hex) < 0.6
        self._slide_fg = design.knockout_text(field_hex, self.colors)
        if not self._editorial:
            return
        try:
            inches = self._Inches
            wash = self._rect(
                slide, inches(0), inches(0), inches(_SLIDE_W_IN), inches(_SLIDE_H_IN), field_hex
            )
            wash.shadow.inherit = False
            self._apply_gradient(wash, [(0.0, field_hex), (1.0, design.darken(field_hex, 16))])
        except Exception:  # noqa: BLE001 — depth is decorative (REQ-5)
            pass

    def _display_headline(
        self,
        slide: Any,
        text: str,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        size: int,
        base_color: str,
        accent_color: str,
        serif_token: bool = False,
        anchor: Any = None,
    ) -> Any:
        """Two-tone display headline; the accent token may render in the serif face (multi-font)."""
        inches = self._Inches
        before, token, after = design.split_for_highlight(text)
        box = slide.shapes.add_textbox(inches(left), inches(top), inches(width), inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        if anchor is not None:
            frame.vertical_anchor = anchor
        para = frame.paragraphs[0]
        para.alignment = self._PP_ALIGN.LEFT

        def _run(segment: str, color: str, *, serif: bool = False) -> None:
            run = para.add_run()
            run.text = segment
            run.font.size = self._Pt(size)
            run.font.bold = True
            run.font.name = self._font_serif if serif else self._font_display
            run.font.color.rgb = self.rgb(color)

        if token:
            if before:
                _run(before, base_color)
            _run(token, accent_color, serif=serif_token)
            if after:
                _run(after, base_color)
        else:
            _run(text, base_color)
        return box

    def _big_number(
        self,
        slide: Any,
        token: str,
        label: str,
        *,
        left: float,
        top: float,
        width: float,
        color: str,
    ) -> None:
        """A hero numeral (grotesk display) with a small mono-ish label beneath it."""
        inches = self._Inches
        self._text(
            slide,
            token,
            inches(left),
            inches(top),
            inches(width),
            inches(1.5),
            92,
            color,
            bold=True,
        )
        if label:
            self._text(
                slide,
                self._short(label, 60),
                inches(left),
                inches(top + 1.45),
                inches(width),
                inches(0.7),
                14,
                self._slide_fg,
                font=self._font_body,
            )

    def _color_card(
        self,
        slide: Any,
        line: str,
        *,
        left: float,
        y: float,
        width: float,
        height: float,
        field: str,
        index: int,
    ) -> None:
        """A saturated color-block card: field fill, serif knockout numeral, body, arrow."""
        inches = self._Inches
        knock = design.knockout_text(field, self.colors)
        card = self._rounded(slide, inches(left), inches(y), inches(width), inches(height), field)
        self._soft_shadow(card)
        self._text(
            slide,
            f"{index:02d}",
            inches(left + 0.32),
            inches(y + 0.2),
            inches(width - 0.6),
            inches(0.8),
            34,
            knock,
            bold=True,
            font=self._font_serif,
        )
        self._text(
            slide,
            "↗",
            inches(left + width - 0.5),
            inches(y + 0.18),
            inches(0.4),
            inches(0.32),
            13,
            knock,
            bold=True,
            align=self._PP_ALIGN.RIGHT,
            font=self._font_body,
        )
        body = self._text(
            slide,
            self._short(line, 150),
            inches(left + 0.34),
            inches(y + 1.05),
            inches(width - 0.66),
            inches(max(0.5, height - 1.42)),
            13,
            knock,
            font=self._font_body,
        )
        body.text_frame.auto_size = self._MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # === Block 2: inch-based surface primitives (consumed by core/blocks via ``Surface``) =====
    def _align(self, name: str) -> Any:
        """Map an alignment name to the pptx enum (blocks pass strings, never pptx enums)."""
        return {
            "left": self._PP_ALIGN.LEFT,
            "center": self._PP_ALIGN.CENTER,
            "right": self._PP_ALIGN.RIGHT,
        }.get(name, self._PP_ALIGN.LEFT)

    def _anchor_of(self, name: str | None) -> Any:
        """Map a vertical-anchor name to the pptx enum (or ``None`` to leave it unset)."""
        if name is None:
            return None
        return {
            "top": self._MSO_ANCHOR.TOP,
            "middle": self._MSO_ANCHOR.MIDDLE,
            "bottom": self._MSO_ANCHOR.BOTTOM,
        }.get(name)

    def fill_region(
        self,
        slide: Any,
        region: layout.Region,
        fill: str,
        *,
        rounded: bool = False,
        line: str | None = None,
        shadow: bool = False,
    ) -> Any:
        """Draw a (rounded) rectangle filling ``region``; optional soft shadow (editorial)."""
        inches = self._Inches
        draw = self._rounded if rounded else self._rect
        shape = draw(
            slide,
            inches(region.left),
            inches(region.top),
            inches(region.width),
            inches(region.height),
            fill,
            line=line,
        )
        if shadow:
            self._soft_shadow(shape)
        return shape

    def text_region(
        self,
        slide: Any,
        region: layout.Region,
        text: str,
        *,
        size: float,
        color: str,
        font: str | None = None,
        bold: bool = False,
        italic: bool = False,
        align: str = "left",
        anchor: str | None = None,
        autofit: bool = False,
        wrap: bool = True,
    ) -> Any:
        """A single-paragraph text box placed in ``region`` (inch-based; used by blocks)."""
        inches = self._Inches
        box = slide.shapes.add_textbox(
            inches(region.left), inches(region.top), inches(region.width), inches(region.height)
        )
        frame = box.text_frame
        frame.word_wrap = wrap
        anc = self._anchor_of(anchor)
        if anc is not None:
            frame.vertical_anchor = anc
        para = frame.paragraphs[0]
        para.text = text
        para.alignment = self._align(align)
        run = para.runs[0] if para.runs else para.add_run()
        run.font.size = self._Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font or (self._font_display if size >= 24 else self._font_body)
        run.font.color.rgb = self.rgb(color)
        if autofit:
            frame.auto_size = self._MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        return box

    def runs_region(
        self,
        slide: Any,
        region: layout.Region,
        runs: Any,
        *,
        align: str = "left",
        anchor: str | None = None,
        wrap: bool = True,
    ) -> Any:
        """A multi-run text box (mixed face/weight/italic/color) placed in ``region``."""
        inches = self._Inches
        box = slide.shapes.add_textbox(
            inches(region.left), inches(region.top), inches(region.width), inches(region.height)
        )
        frame = box.text_frame
        frame.word_wrap = wrap
        anc = self._anchor_of(anchor)
        if anc is not None:
            frame.vertical_anchor = anc
        para = frame.paragraphs[0]
        para.alignment = self._align(align)
        for r in runs:
            run = para.add_run()
            run.text = r.text
            run.font.size = self._Pt(r.size)
            run.font.bold = r.bold
            run.font.italic = r.italic
            run.font.name = r.font
            run.font.color.rgb = self.rgb(r.color)
        return box

    def image_region(
        self,
        slide: Any,
        region: layout.Region,
        path: str,
        *,
        cover: bool = True,
        scrim_alpha: int | None = None,
    ) -> bool:
        """Place an image cover-fit (cropped, no stretch) into ``region``; optional dark scrim."""
        inches = self._Inches
        try:
            pic = slide.shapes.add_picture(str(path), inches(region.left), inches(region.top))
            if cover and pic.width and pic.height:
                aspect = pic.width / pic.height
                target = region.width / region.height
                if aspect > target:
                    crop = (1 - target / aspect) / 2
                    pic.crop_left = crop
                    pic.crop_right = crop
                else:
                    crop = (1 - aspect / target) / 2
                    pic.crop_top = crop
                    pic.crop_bottom = crop
            pic.left = inches(region.left)
            pic.top = inches(region.top)
            pic.width = inches(region.width)
            pic.height = inches(region.height)
            if scrim_alpha is not None:
                scrim = self._rect(
                    slide,
                    inches(region.left),
                    inches(region.top),
                    inches(region.width),
                    inches(region.height),
                    self.colors.canvas_dark,
                )
                scrim.shadow.inherit = False
                self._set_alpha(scrim, scrim_alpha)
            return True
        except Exception:  # noqa: BLE001 — a bad image never breaks a slide (REQ-5)
            return False

    def gradient(self, shape: Any, stops: list[tuple[float, str]]) -> None:
        """Public alias: replace ``shape``'s fill with a multi-stop gradient (depth)."""
        self._apply_gradient(shape, stops)

    def set_alpha(self, shape: Any, pct: int) -> None:
        """Public alias: set a solid-filled shape's opacity (0..100)."""
        self._set_alpha(shape, pct)

    def soft_shadow(self, shape: Any) -> None:
        """Public alias: add the editorial soft outer shadow (fail-open)."""
        self._soft_shadow(shape)

    def card_system(
        self,
        slide: Any,
        line: str,
        *,
        left: float,
        y: float,
        width: float,
        height: float,
        accent: str,
        index: int,
    ) -> None:
        """Public alias for the white system card (accent spine, metric/index, arrow, tag)."""
        self._system_card(
            slide, line, left=left, y=y, width=width, height=height, accent=accent, index=index
        )

    def card_color(
        self,
        slide: Any,
        line: str,
        *,
        left: float,
        y: float,
        width: float,
        height: float,
        field: str,
        index: int,
    ) -> None:
        """Public alias for the saturated color-block card (knockout numeral, body, arrow)."""
        self._color_card(
            slide, line, left=left, y=y, width=width, height=height, field=field, index=index
        )

    def big_number(
        self,
        slide: Any,
        token: str,
        label: str,
        *,
        left: float,
        top: float,
        width: float,
        color: str,
    ) -> None:
        """Public alias for the hero numeral + label."""
        self._big_number(slide, token, label, left=left, top=top, width=width, color=color)

    def short(self, text: str, limit: int = 165) -> str:
        """Public alias for sentence/word-boundary truncation (the 4B-safety trim)."""
        return self._short(text, limit)

    def metric_token(self, text: str) -> str | None:
        """Public alias for the visible-number-token extractor."""
        return self._metric_token(text)

    def _new(self) -> Any:
        """Append a blank slide."""
        return self.prs.slides.add_slide(self._blank)

    def _try_chart(self, slide: Any, sc: SlideContent, *, top: float) -> bool:
        """Render the slide's ``.visual`` as a themed image-chart (fail-open to a native chart).

        Respects ``style.max_charts_per_deck`` and the master ``charts_enabled`` switch. Prefers the
        magazine-grade matplotlib image-chart (Editorial palette + the deck's fonts, labeled axes);
        if that fails it falls back to the native editable chart, then to cards/table. Returns False
        when there is no visual, the budget is spent, charts are disabled, or both renderers fail.
        """
        if sc.visual is None or not self.style.charts_enabled:
            return False
        if self._charts_used >= self._max_charts:
            return False
        height = max(2.4, min(3.7, 6.55 - top))
        ok = charts_image.add_image_chart(
            slide,
            sc.visual,
            self.colors,
            self.style,
            left_in=0.72,
            top_in=top,
            width_in=11.75,
            height_in=height,
            on_dark=self._slide_on_dark,
        )
        if not ok:
            ok = visuals.add_native_chart(
                slide,
                sc.visual,
                self.colors,
                left_in=0.72,
                top_in=top,
                width_in=11.75,
                height_in=height,
            )
        if ok:
            self._charts_used += 1
        return ok

    # --- native diagrams (Schaubilder) ---------------------------------------------------
    def _try_diagram(self, slide: Any, sc: SlideContent, *, top: float) -> bool:
        """Render the slide's ``.diagram`` as native shapes (fail-open; respects the budget)."""
        if sc.diagram is None or not self.style.charts_enabled:
            return False
        if self._diagrams_used >= self._max_diagrams:
            return False
        try:
            ok = self._render_diagram(slide, sc.diagram, top=top)
        except Exception:  # noqa: BLE001 — a diagram quirk never loses the slide (REQ-5)
            ok = False
        if ok:
            self._diagrams_used += 1
        return ok

    def _render_diagram(self, slide: Any, spec: DiagramSpec, *, top: float) -> bool:
        """Dispatch a :class:`DiagramSpec` to its native shape renderer."""
        kind = spec.diagram_type
        if kind == DiagramType.PROCESS:
            return self._diagram_process(slide, spec, top=top)
        if kind == DiagramType.MATRIX_2X2:
            return self._diagram_matrix(slide, spec, top=top)
        if kind in (DiagramType.PYRAMID, DiagramType.FUNNEL):
            return self._diagram_stack(slide, spec, top=top, funnel=kind == DiagramType.FUNNEL)
        if kind == DiagramType.KPI_STRIP:
            return self._diagram_kpi(slide, spec, top=top)
        if kind == DiagramType.COMPARE_COLUMNS:
            return self._diagram_columns(slide, spec, top=top)
        return False

    def _diagram_process(self, slide: Any, spec: DiagramSpec, *, top: float) -> bool:
        """Ordered process flow: saturated cards left→right joined by arrows."""
        inches = self._Inches
        fields = design.statement_fields(self.colors)
        n = len(spec.nodes)
        gap = 0.3
        card_w = (11.75 - gap * (n - 1)) / n
        height = 2.0
        for i, node in enumerate(spec.nodes):
            left = 0.72 + i * (card_w + gap)
            self._color_card(
                slide,
                node.label,
                left=left,
                y=top,
                width=card_w,
                height=height,
                field=fields[i % len(fields)],
                index=i + 1,
            )
            if i < n - 1:
                self._text(
                    slide,
                    "→",
                    inches(left + card_w - 0.05),
                    inches(top + height / 2 - 0.2),
                    inches(gap + 0.1),
                    inches(0.4),
                    20,
                    self._slide_fg,
                    bold=True,
                    align=self._PP_ALIGN.CENTER,
                )
        return True

    def _diagram_matrix(self, slide: Any, spec: DiagramSpec, *, top: float) -> bool:
        """2x2 matrix: four saturated quadrants with knockout label + detail."""
        inches = self._Inches
        fields = design.statement_fields(self.colors)
        positions = [(0.72, top), (6.78, top), (0.72, top + 2.45), (6.78, top + 2.45)]
        for i, (node, (left, y)) in enumerate(zip(spec.nodes, positions, strict=False)):
            field = fields[i % len(fields)]
            knock = design.knockout_text(field, self.colors)
            card = self._rounded(slide, inches(left), inches(y), inches(5.68), inches(2.3), field)
            self._soft_shadow(card)
            self._text(
                slide,
                node.label,
                inches(left + 0.25),
                inches(y + 0.16),
                inches(5.2),
                inches(0.5),
                16,
                knock,
                bold=True,
            )
            if node.detail:
                self._text(
                    slide,
                    self._short(node.detail, 120),
                    inches(left + 0.25),
                    inches(y + 0.72),
                    inches(5.2),
                    inches(1.45),
                    12,
                    knock,
                    font=self._font_body,
                )
        return True

    def _diagram_stack(self, slide: Any, spec: DiagramSpec, *, top: float, funnel: bool) -> bool:
        """Pyramid (broad base) or funnel (narrowing) of centered saturated tiers."""
        inches = self._Inches
        fields = design.statement_fields(self.colors)
        n = len(spec.nodes)
        height = min(0.95, 4.6 / n)
        gap = 0.14
        for i, node in enumerate(spec.nodes):
            t = i / (n - 1) if n > 1 else 0.0
            wfrac = (1.0 - 0.6 * t) if funnel else (0.45 + 0.55 * t)
            width = 11.75 * wfrac
            left = 0.72 + (11.75 - width) / 2
            y = top + i * (height + gap)
            field = fields[i % len(fields)]
            knock = design.knockout_text(field, self.colors)
            bar = self._rounded(
                slide, inches(left), inches(y), inches(width), inches(height), field
            )
            self._soft_shadow(bar)
            label = node.label + (f"  ·  {node.value}" if node.value else "")
            self._text(
                slide,
                self._short(label, 70),
                inches(left + 0.2),
                inches(y),
                inches(width - 0.4),
                inches(height),
                13,
                knock,
                bold=True,
                anchor=self._MSO_ANCHOR.MIDDLE,
                align=self._PP_ALIGN.CENTER,
            )
        return True

    def _diagram_kpi(self, slide: Any, spec: DiagramSpec, *, top: float) -> bool:
        """KPI strip: 2–5 saturated tiles each with a big value over a label."""
        inches = self._Inches
        fields = design.statement_fields(self.colors)
        n = len(spec.nodes)
        gap = 0.3
        width = (11.75 - gap * (n - 1)) / n
        height = 2.2
        for i, node in enumerate(spec.nodes):
            left = 0.72 + i * (width + gap)
            field = fields[i % len(fields)]
            knock = design.knockout_text(field, self.colors)
            tile = self._rounded(
                slide, inches(left), inches(top), inches(width), inches(height), field
            )
            self._soft_shadow(tile)
            self._text(
                slide,
                self._short(node.value, 10),
                inches(left + 0.15),
                inches(top + 0.28),
                inches(width - 0.3),
                inches(1.0),
                40,
                knock,
                bold=True,
                align=self._PP_ALIGN.CENTER,
                anchor=self._MSO_ANCHOR.MIDDLE,
            )
            self._text(
                slide,
                self._short(node.label, 24),
                inches(left + 0.15),
                inches(top + 1.4),
                inches(width - 0.3),
                inches(0.7),
                12,
                knock,
                font=self._font_body,
                align=self._PP_ALIGN.CENTER,
            )
        return True

    def _diagram_columns(self, slide: Any, spec: DiagramSpec, *, top: float) -> bool:
        """Comparison columns: 2–3 saturated panels with a serif title + attribute rows."""
        inches = self._Inches
        fields = design.statement_fields(self.colors)
        m = len(spec.columns)
        gap = 0.3
        width = (11.75 - gap * (m - 1)) / m
        height = min(4.6, 6.55 - top)
        for i, column in enumerate(spec.columns):
            left = 0.72 + i * (width + gap)
            field = fields[i % len(fields)]
            knock = design.knockout_text(field, self.colors)
            panel = self._rounded(
                slide, inches(left), inches(top), inches(width), inches(height), field
            )
            self._soft_shadow(panel)
            self._text(
                slide,
                self._short(column.title, 18),
                inches(left + 0.2),
                inches(top + 0.18),
                inches(width - 0.4),
                inches(0.6),
                18,
                knock,
                bold=True,
                font=self._font_serif,
            )
            cy = top + 0.95
            for cell in column.cells[:5]:
                self._text(
                    slide,
                    self._short(cell, 42),
                    inches(left + 0.2),
                    inches(cy),
                    inches(width - 0.4),
                    inches(0.6),
                    12,
                    knock,
                    font=self._font_body,
                )
                cy += 0.62
        return True

    # --- per-slide-type renderers --------------------------------------------------------
    def _title_slide(self, sc: SlideContent) -> None:
        """Render the cover: a luminous dark editorial canvas + a two-tone grotesk headline."""
        inches = self._Inches
        slide = self._new()
        self._slide_no += 1
        self._slide_on_dark = True
        self._slide_fg = design.contrast_text(self.colors.canvas_dark, self.colors)
        cover = imagery.cover_image(self._imagery_dir, seed=sc.headline)
        if cover is not None:
            self._paint_cover_image(slide, cover)  # full-bleed Neura image under a dark scrim
        else:
            self._paint_dark_canvas(slide, glow=True)
        full = inches(_SLIDE_H_IN)
        self._rect(slide, inches(0), inches(0), inches(0.22), full, self.colors.accent_cyan)
        self._rect(slide, inches(0.22), inches(0), inches(0.08), full, self.colors.artifact_gold)
        self._text(
            slide,
            framework_marker().upper(),
            inches(0.82),
            inches(0.72),
            inches(8.0),
            inches(0.24),
            9,
            self.colors.accent_cyan,
            bold=True,
            font=self._font_mono,
        )
        self._headline_two_tone(
            slide,
            sc.headline,
            left=0.8,
            top=2.45,
            width=11.7,
            height=2.7,
            size=40,
            base_color=self._slide_fg,
            accent_color=self.colors.artifact_gold,
        )
        subtitle = sc.body or date.today().isoformat()
        self._text(
            slide,
            subtitle,
            inches(0.82),
            inches(5.25),
            inches(11.6),
            inches(0.8),
            18,
            self.colors.accent_cyan,
        )
        self._text(
            slide,
            date.today().isoformat(),
            inches(10.2),
            inches(0.7),
            inches(2.1),
            inches(0.3),
            10,
            self.colors.light_surface,
            bold=True,
            align=self._PP_ALIGN.RIGHT,
            font=self._font_mono,
        )
        self._add_logo(slide)

    def _table_slide(self, sc: SlideContent, canvas: str | None = None) -> None:
        slide = self._new()
        self._frame(slide, sc, canvas_hex=canvas)
        self._headline(slide, sc.headline, accent=self._accent(sc.slide_type))
        top = 1.75
        if sc.body:
            self._body_callout(slide, sc.body, top=1.58, fill=self.colors.deep_blue)
            top = 3.05
        if sc.table:
            self._draw_table(slide, sc.table, top=top)
        elif self._try_chart(slide, sc, top=top):
            return
        elif self._try_diagram(slide, sc, top=top):
            return
        else:
            self._signal_cards(slide, sc.bullets, top=top)

    def _draw_table(self, slide: Any, rows: list[list[str]], *, top: float = 1.8) -> None:
        """Draw a header-styled table (row 0 = header) with alternating row fills."""
        inches, pt = self._Inches, self._Pt
        n_rows = len(rows)
        n_cols = max((len(r) for r in rows), default=1)
        height = min(4.9, max(0.7, 0.48 * n_rows))
        graphic = slide.shapes.add_table(
            n_rows, n_cols, inches(0.72), inches(top), inches(11.75), inches(height)
        )
        table = graphic.table
        for r, row in enumerate(rows):
            for c in range(n_cols):
                cell = table.cell(r, c)
                cell.text = self._short(str(row[c]) if c < len(row) else "", 95)
                para = cell.text_frame.paragraphs[0]
                run = para.runs[0] if para.runs else para.add_run()
                run.font.name = self._font_body
                run.font.size = pt(12 if r else 13)
                run.font.bold = r == 0
                if r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.rgb(self.colors.excel_header)
                    run.font.color.rgb = self.rgb(self.colors.white)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.rgb(
                        self.colors.white if r % 2 else self.colors.light_surface
                    )
                    run.font.color.rgb = self.rgb(self.colors.text_dark)

    def _matrix_slide(self, sc: SlideContent, canvas: str | None = None) -> None:
        """Render a 2x2 matrix (SWOT / positioning) as four saturated knockout quadrants."""
        slide = self._new()
        self._frame(slide, sc, canvas_hex=canvas)
        self._headline(slide, sc.headline, accent=self._accent(sc.slide_type))
        quadrants = self._swot_quadrants(sc)
        nodes = [
            DiagramNode(label=label, detail="; ".join(lines[:3])) for label, lines in quadrants
        ]
        try:
            spec = DiagramSpec(diagram_type=DiagramType.MATRIX_2X2, nodes=nodes)
        except ValueError:
            self._signal_cards(slide, [lbl for lbl, _ in quadrants], top=1.95)
            return
        self._diagram_matrix(slide, spec, top=1.95)

    def _swot_slide(self, sc: SlideContent, canvas: str | None = None) -> None:
        """Back-compatible SWOT entry — delegates to the generalized 2x2 matrix renderer."""
        self._matrix_slide(sc, canvas)

    def _swot_quadrants(self, sc: SlideContent) -> list[tuple[str, list[str]]]:
        """Resolve 4 (label, lines) quadrants from the slide's table or bullets (fail-safe)."""
        default_labels = ["Strengths", "Weaknesses", "Opportunities", "Threats"]
        if sc.table and len(sc.table) >= 4:
            out: list[tuple[str, list[str]]] = []
            for row in sc.table[:4]:
                label = str(row[0]) if row else ""
                body = str(row[1]) if len(row) > 1 else ""
                out.append(
                    (label, [seg.strip() for seg in re.split(r"[;\n]", body) if seg.strip()])
                )
            return out
        bullets = sc.bullets or []
        return [(default_labels[i], [bullets[i]] if i < len(bullets) else []) for i in range(4)]

    def _statement_slide(self, sc: SlideContent, canvas: str) -> None:
        """Full-bleed manifesto: a saturated field, an oversized two-tone serif-accent headline."""
        inches = self._Inches
        slide = self._new()
        self._slide_no += 1
        self._paint_field(slide, canvas)
        fg = self._slide_fg
        spot = design.spot_for_canvas(canvas, self.colors)
        self._text(
            slide,
            deck_frame_label(self.language, sc.slide_type).upper(),
            inches(0.82),
            inches(0.7),
            inches(10.0),
            inches(0.3),
            11,
            fg,
            bold=True,
            font=self._font_mono,
        )
        self._display_headline(
            slide,
            sc.headline,
            left=0.8,
            top=1.9,
            width=11.7,
            height=3.4,
            size=56,
            base_color=fg,
            accent_color=spot,
            serif_token=True,
        )
        if sc.body:
            self._text(
                slide,
                self._short(sc.body, 120),
                inches(0.82),
                inches(5.7),
                inches(10.2),
                inches(1.0),
                18,
                fg,
                font=self._font_body,
            )
        self._text(
            slide,
            f"{self._slide_no:02d}",
            inches(10.9),
            inches(5.5),
            inches(2.1),
            inches(1.6),
            80,
            spot,
            bold=True,
            align=self._PP_ALIGN.RIGHT,
        )
        self._add_logo(slide)

    def _quote_slide(self, sc: SlideContent, canvas: str) -> None:
        """Oversized serif pull-statement on a dark/saturated canvas with a quote glyph."""
        inches = self._Inches
        slide = self._new()
        self._slide_no += 1
        if canvas == self.colors.canvas_dark:
            self._slide_on_dark = True
            self._slide_fg = design.contrast_text(canvas, self.colors)
            self._paint_dark_canvas(slide, glow=True)
        else:
            self._paint_field(slide, canvas)
        fg = self._slide_fg
        spot = design.spot_for_canvas(canvas, self.colors)
        self._text(
            slide,
            "“",
            inches(0.55),
            inches(0.35),
            inches(3.0),
            inches(2.4),
            130,
            spot,
            bold=True,
            font=self._font_serif,
        )
        quote = sc.body or sc.headline
        self._text(
            slide,
            self._short(quote, 170),
            inches(1.2),
            inches(2.1),
            inches(11.0),
            inches(3.2),
            40,
            fg,
            font=self._font_serif,
        )
        if sc.body and sc.headline:
            self._text(
                slide,
                self._short(sc.headline, 40),
                inches(1.2),
                inches(5.7),
                inches(10.0),
                inches(0.5),
                13,
                fg,
                font=self._font_mono,
            )
        self._add_logo(slide)

    def _metric_hero_slide(self, sc: SlideContent, canvas: str | None = None) -> None:
        """One–three giant grounded numerals with a supporting clause beneath each."""
        slide = self._new()
        self._frame(slide, sc, canvas_hex=canvas)
        self._headline(slide, sc.headline, accent=self._accent(sc.slide_type))
        metrics: list[tuple[str, str]] = []
        for bullet in sc.bullets:
            token = self._metric_token(bullet)
            if token:
                metrics.append((token, bullet))
        metrics = metrics[:3]
        if not metrics:
            self._signal_cards(slide, sc.bullets, top=1.95)
            return
        gap = 0.4
        width = (11.75 - gap * (len(metrics) - 1)) / len(metrics)
        spot = self._spot_color()
        for i, (token, line) in enumerate(metrics):
            left = 0.72 + i * (width + gap)
            self._big_number(slide, token, line, left=left, top=2.2, width=width, color=spot)

    def _colorblock_grid_slide(
        self, sc: SlideContent, canvas: str | None = None, accent_index: int = 0
    ) -> None:
        """Saturated numbered color-block cards (the KINETIC / 'Selected Work' look)."""
        slide = self._new()
        self._frame(slide, sc, canvas_hex=canvas)
        self._headline(slide, sc.headline, accent=self._accent(sc.slide_type))
        lines = [self._short(b, 150) for b in (sc.bullets or []) if str(b).strip()][:4]
        if not lines:
            self._signal_cards(slide, sc.bullets, top=1.95)
            return
        top = 1.95
        if sc.body:
            self._body_callout(slide, sc.body, top=1.62, fill=self.colors.deep_blue)
            top = 3.05
        fields = design.statement_fields(self.colors)
        count = len(lines)
        cols = 1 if count == 1 else 2
        rows = (count + cols - 1) // cols
        gap = 0.3
        card_w = 11.75 if cols == 1 else 5.73
        available = max(1.1, 6.55 - top)
        card_h = min(2.3, (available - gap * (rows - 1)) / rows)
        for idx, line in enumerate(lines):
            row, col = divmod(idx, cols)
            left = 0.72 + col * (card_w + gap)
            y = top + row * (card_h + gap)
            field = fields[(accent_index + idx) % len(fields)]
            self._color_card(
                slide, line, left=left, y=y, width=card_w, height=card_h, field=field, index=idx + 1
            )

    def _editorial_split_slide(self, sc: SlideContent, canvas: str | None = None) -> None:
        """Asymmetric editorial split: a big serif headline left, supporting body/bullets right."""
        inches = self._Inches
        slide = self._new()
        self._frame(slide, sc, canvas_hex=canvas)
        spot = self._spot_color()
        self._display_headline(
            slide,
            sc.headline,
            left=0.7,
            top=1.9,
            width=6.6,
            height=3.9,
            size=44,
            base_color=self._slide_fg,
            accent_color=spot,
            serif_token=True,
        )
        self._rect(
            slide, inches(7.55), inches(1.95), inches(0.012), inches(3.7), self.colors.charcoal
        )
        rx, ry = 7.95, 1.98
        if sc.body:
            self._text(
                slide,
                self._short(sc.body, 220),
                inches(rx),
                inches(ry),
                inches(4.65),
                inches(1.7),
                14,
                self._slide_fg,
                font=self._font_body,
            )
            ry += 1.75
        for bullet in (sc.bullets or [])[:4]:
            self._text(
                slide,
                "— " + self._short(bullet, 90),
                inches(rx),
                inches(ry),
                inches(4.65),
                inches(0.7),
                13,
                self._slide_fg,
                font=self._font_body,
            )
            ry += 0.72

    def _chart_slide(self, sc: SlideContent, canvas: str | None = None) -> None:
        """A data slide: native chart (or native diagram), falling back to table/cards."""
        slide = self._new()
        self._frame(slide, sc, canvas_hex=canvas)
        self._headline(slide, sc.headline, accent=self._accent(sc.slide_type))
        top = 1.85
        if sc.body:
            self._body_callout(slide, sc.body, top=1.62, fill=self._accent(sc.slide_type))
            top = 3.05
        if self._try_chart(slide, sc, top=top):
            return
        if self._try_diagram(slide, sc, top=top):
            return
        if sc.table:
            self._draw_table(slide, sc.table, top=top)
        else:
            self._signal_cards(slide, sc.bullets, top=top)

    def _appendix_slide(self, sc: SlideContent, canvas: str | None = None) -> None:
        """Render sources/appendix material as a compact two-column reference page."""
        slide = self._new()
        self._frame(slide, sc, canvas_hex=canvas)
        self._headline(slide, sc.headline, accent=self._accent(sc.slide_type))
        if sc.body:
            self._body_callout(slide, sc.body, top=1.55, fill=self.colors.charcoal)
            top = 3.0
        else:
            top = 1.72
        items = sc.bullets
        if sc.table:
            items = [" - ".join(str(cell) for cell in row if str(cell).strip()) for row in sc.table]
        self._compact_list(slide, items, top=top)

    def _content_slide(self, sc: SlideContent, canvas: str | None = None) -> None:
        """Render a content slide through the frame (chart, diagram, table, or cards)."""
        slide = self._new()
        self._frame(slide, sc, canvas_hex=canvas)
        self._headline(slide, sc.headline, accent=self._accent(sc.slide_type))
        top = 1.75
        if sc.body:
            self._body_callout(slide, sc.body, top=1.58, fill=self._accent(sc.slide_type))
            top = 3.05
        if self._try_chart(slide, sc, top=top):
            return
        if self._try_diagram(slide, sc, top=top):
            return
        if sc.table:
            self._draw_table(slide, sc.table, top=top)
        else:
            self._signal_cards(slide, sc.bullets, top=top)

    def render(
        self,
        sc: SlideContent,
        plan: SlidePlan | None = None,
        comp: composer.SlideComposition | None = None,
    ) -> None:
        """Render a slide: the Block-2 :class:`SlideComposition` first, the archetype path next.

        When ``comp`` is supplied (the default from :func:`build_deck`) the composable library
        paints the slide. ``_render_composition`` never raises (each block is fail-open), so the
        only way to reach the legacy archetype path is a missing composition — preserving REQ-5's
        ultimate fallback (and the pre-Block-2 behavior for callers that pass only a ``plan``).
        """
        if comp is not None and self._render_composition(sc, comp):
            return
        if sc.slide_type == SlideType.TITLE:
            self._title_slide(sc)
            return
        if plan is None:
            self._content_slide(sc)
            return
        canvas = deck_director.canvas_hex(plan.canvas, self.colors, plan.accent_index)
        try:
            self._render_archetype(plan.archetype, sc, canvas, plan.accent_index)
        except Exception:  # noqa: BLE001 — any archetype quirk degrades to content (REQ-5)
            self._content_slide(sc, canvas)

    # --- Block 2: render a composer SlideComposition -------------------------------------
    def _render_composition(self, sc: SlideContent, comp: composer.SlideComposition) -> bool:
        """Paint a :class:`~core.composer.SlideComposition`; False only if no slide was made.

        Self-contained and fail-open: the canvas, chrome and every block are individually guarded,
        so this never raises — a partial block failure still yields a clean slide with its headline
        rather than handing back to the archetype path (which would add a duplicate slide).
        """
        try:
            slide = self._new()
        except Exception:  # noqa: BLE001 — pptx could not add a slide; let the caller fall back
            return False
        self._slide_no = comp.position
        try:
            self._paint_composition_canvas(slide, comp.canvas)
        except Exception:  # noqa: BLE001 — degrade to a clean cream canvas (REQ-5)
            self._paint_solid_or_dark(slide, self.colors.paper)
        if comp.chrome:
            try:
                self._frame_chrome(slide, sc, number=comp.position)
            except Exception:  # noqa: BLE001 — chrome is decorative; never lose the slide
                pass
        else:
            self._add_logo(slide)
        theme = self._block_theme(comp)
        for placed in comp.blocks:
            if placed.kind == "chart":
                self._render_chart_block(slide, placed.region, placed.params)
                continue
            blocks.render(placed.kind, self, slide, placed.region, placed.params, theme)
        return True

    def _paint_composition_canvas(self, slide: Any, canvas: templates.CanvasSpec) -> None:
        """Paint a composition's canvas (image / dark / field / solid) + set the per-slide fg."""
        if canvas.role == "image" and canvas.image_path:
            self._slide_on_dark = True
            self._slide_fg = design.contrast_text(self.colors.canvas_dark, self.colors)
            self._paint_cover_image(slide, canvas.image_path)
        elif canvas.role == "dark":
            self._slide_on_dark = True
            self._slide_fg = design.contrast_text(self.colors.canvas_dark, self.colors)
            self._paint_dark_canvas(slide, glow=canvas.glow)
        elif canvas.role == "field" and canvas.field_hex:
            self._paint_field(slide, canvas.field_hex)  # sets fg/on_dark itself
        else:
            background = slide.background
            background.fill.solid()
            background.fill.fore_color.rgb = self.rgb(canvas.fill)
            self._slide_on_dark = canvas.on_dark
            self._slide_fg = design.contrast_text(canvas.fill, self.colors)

    def _block_theme(self, comp: composer.SlideComposition) -> blocks.BlockTheme:
        """Build the canvas-derived :class:`~core.blocks.BlockTheme` handed to every block."""
        canvas = comp.canvas
        base = canvas.field_hex or (self.colors.canvas_dark if self._slide_on_dark else canvas.fill)
        muted = self.colors.light_surface if self._slide_on_dark else self.colors.charcoal
        return blocks.BlockTheme(
            colors=self.colors,
            fonts=self.fonts,
            editorial=self._editorial,
            fg=self._slide_fg,
            on_dark=self._slide_on_dark,
            spot=design.spot_for_canvas(base, self.colors),
            muted=muted,
            accent=comp.accent,
        )

    def _render_chart_block(self, slide: Any, region: layout.Region, params: Any) -> None:
        """Render a composition chart block, enforcing the per-deck chart budget (fail-open)."""
        spec = params.get("spec")
        if not isinstance(spec, ChartSpec) or not self.style.charts_enabled:
            return
        if self._charts_used >= self._max_charts:
            return
        height = min(region.height, 3.9)
        ok = charts_image.add_image_chart(
            slide,
            spec,
            self.colors,
            self.style,
            left_in=region.left,
            top_in=region.top,
            width_in=region.width,
            height_in=height,
            on_dark=self._slide_on_dark,
        )
        if not ok:
            ok = visuals.add_native_chart(
                slide,
                spec,
                self.colors,
                left_in=region.left,
                top_in=region.top,
                width_in=region.width,
                height_in=height,
            )
        if ok:
            self._charts_used += 1

    def _render_archetype(
        self, archetype: Archetype, sc: SlideContent, canvas: str, accent_index: int
    ) -> None:
        """Dispatch one slide to its archetype renderer (see :class:`~models.deck.Archetype`)."""
        if archetype == Archetype.STATEMENT:
            self._statement_slide(sc, canvas)
        elif archetype == Archetype.QUOTE:
            self._quote_slide(sc, canvas)
        elif archetype == Archetype.METRIC_HERO:
            self._metric_hero_slide(sc, canvas)
        elif archetype == Archetype.COLORBLOCK_GRID:
            self._colorblock_grid_slide(sc, canvas, accent_index)
        elif archetype == Archetype.EDITORIAL_SPLIT:
            self._editorial_split_slide(sc, canvas)
        elif archetype == Archetype.TABLE:
            self._table_slide(sc, canvas)
        elif archetype == Archetype.MATRIX:
            self._matrix_slide(sc, canvas)
        elif archetype == Archetype.CHART:
            self._chart_slide(sc, canvas)
        elif archetype == Archetype.APPENDIX:
            self._appendix_slide(sc, canvas)
        else:
            self._content_slide(sc, canvas)

    def save(self, output_dir: str | Path, title: str) -> Path:
        """Write the deck to ``output_dir``, embed its fonts for portability, return its path."""
        out = Path(output_dir)
        out.mkdir(parents=True, exist_ok=True)
        path = out / f"{date.today().isoformat()}_{_slug(title)}_deck.pptx"
        self.prs.save(str(path))
        # Embed the deck's font families so it renders correctly when forwarded to a machine that
        # does not have them installed (the CEO/board) — fail-open, never blocks the render (REQ-5).
        font_embed.embed_fonts(
            path,
            [
                self._font_display,
                self._font_body,
                self._font_mono,
                self._font_serif,
                self._font_statement,
            ],
            self.style.fonts_dir,
        )
        return path


def build_deck(
    deck: DeckStructure,
    config: AppConfig,
    output_dir: str | Path,
    analysis: AnalysisOutput | None = None,
    research_report: ResearchReport | None = None,
) -> Path:
    """Render a :class:`~models.deck.DeckStructure` into a Porter Editorial .pptx; return its path.

    Handles all 10 SPEC §11 slide types in the Editorial system (cream/dark canvas, multi-font,
    two-tone headlines, system cards, native charts on a slide's ``.visual``, logo bottom-right).
    When ``research_report`` is supplied, source-grounded telemetry chips (DNA 6) render on the
    bottom rail — the pipeline threads it in (Block 4). Raises :class:`DeckBuildError` if
    python-pptx is unavailable.
    """
    deck = prepare_deck_for_render(deck, analysis)
    renderer = _DeckRenderer(config, deck.language, research_report)
    editorial = design.is_editorial(config.output.style)
    ctx = composer.DeckContext(
        colors=config.output.colors,
        style=config.output.style,
        language=deck.language,
        editorial=editorial,
        imagery_dir=config.output.imagery_dir,
    )
    comps = composer.compose_deck(deck.slides, ctx)
    plans = deck_director.plan_deck(deck.slides, editorial=editorial)  # ultimate-fallback path
    for slide, plan, comp in zip(deck.slides, plans, comps, strict=True):
        renderer.render(slide, plan, comp)
    return renderer.save(output_dir, deck.title)


def management_deck_structure(analysis: AnalysisOutput, language: Language) -> DeckStructure:
    """Build a generic management :class:`DeckStructure` from an analysis (doc-prep / fallback).

    Title → executive summary → one 'so what' slide per section → sources. This is the
    deterministic mapping used when no LLM deck shaping is applied.
    """
    is_de = language == Language.DE
    slides = [
        SlideContent(
            slide_type=SlideType.TITLE,
            headline=analysis.title,
            body=("Management-Briefing" if is_de else "Management Briefing")
            + f"  ·  {date.today().isoformat()}",
        ),
        SlideContent(
            slide_type=SlideType.EXECUTIVE_SUMMARY,
            headline="Kernaussage" if is_de else "Executive Summary",
            bullets=_sentences(analysis.bottom_line),
        ),
    ]
    for section in analysis.sections:
        slides.append(
            SlideContent(
                slide_type=SlideType.STRATEGIC_SIGNALS,
                headline=section.heading,
                bullets=_sentences(section.body),
            )
        )
    if analysis.sources:
        slides.append(
            SlideContent(
                slide_type=SlideType.APPENDIX,
                headline="Quellen" if is_de else "Sources",
                bullets=[s.url + (f" — {s.title}" if s.title else "") for s in analysis.sources],
            )
        )
    return DeckStructure(title=analysis.title, language=language, slides=slides)


def build_management_deck(
    analysis: AnalysisOutput, config: AppConfig, output_dir: str | Path, language: Language
) -> Path:
    """Back-compatible management deck — builds a generic structure and renders via build_deck."""
    return build_deck(
        management_deck_structure(analysis, language),
        config,
        output_dir,
        analysis=analysis,
    )


# ----------------------------------------------------------------- PDF (back-compat shim)
def build_management_pdf(
    analysis: AnalysisOutput, config: AppConfig, output_dir: str | Path, language: Language
) -> Path:
    """Back-compatible management PDF brief — delegates to :func:`build_brief_pdf` (T-5 layout).

    Kept so the doc-prep path and existing tests keep working; new callers should use
    :func:`build_brief_pdf` directly with the task type.
    """
    return build_brief_pdf(analysis, config, output_dir, task_type=TaskType.DOCUMENT_SYNTHESIS)
